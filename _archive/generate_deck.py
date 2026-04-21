#!/usr/bin/env python3
"""
Generate SLASHR presentation deck (PPTX) for Ahold Delhaize / Carlos Vicente.
Dark-first design, brand colors, professional layout.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Brand colors ---
BG_DARK = RGBColor(0x1A, 0x1A, 0x1A)
SURFACE = RGBColor(0x2C, 0x2E, 0x34)
SURFACE_ALT = RGBColor(0x25, 0x27, 0x2E)
ORANGE = RGBColor(0xE7, 0x46, 0x01)
MAGENTA = RGBColor(0xCE, 0x08, 0xA9)
VIOLET = RGBColor(0x89, 0x62, 0xFD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WHITE_70 = RGBColor(0xB3, 0xB3, 0xB3)  # 70% white on dark
WHITE_50 = RGBColor(0x80, 0x80, 0x80)  # 50% white on dark
ORANGE_LIGHT = RGBColor(0xFF, 0x90, 0x11)

# --- Dimensions ---
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.8)
MARGIN_R = Inches(0.8)
CONTENT_W = Inches(11.733)

LOGO_PATH = "/Users/quentin/Desktop/LOGO-SLASHR-BLANC-1.png"
JANUS_PATH = "/Users/quentin/Downloads/iPhone-13-PRO-MAX-janus.agence-slashr.fr.png"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color=SURFACE, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius is not None:
        # Adjustment in EMU for corner radius
        shape.adjustments[0] = radius
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="Inter", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.line_spacing = Pt(font_size * line_spacing)
    return txBox, tf


def add_paragraph(tf, text, font_size=18, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, font_name="Inter", space_before=0, line_spacing=1.2):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.line_spacing = Pt(font_size * line_spacing)
    return p


def add_accent_bar(slide, left, top, width=Inches(0.6), height=Pt(4), color=ORANGE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_gradient_bar(slide, left, top, width=Inches(2), height=Pt(4)):
    """Simulated gradient bar with 3 segments"""
    seg_w = width // 3
    for i, color in enumerate([ORANGE, MAGENTA, VIOLET]):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + seg_w * i, top, seg_w, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()


def add_tag(slide, left, top, text, color=WHITE_50, bg_color=SURFACE):
    w, h = Inches(3.5), Pt(28)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = "Inter"
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_logo(slide, left=None, top=None, width=Inches(1.5)):
    if left is None:
        left = MARGIN_L
    if top is None:
        top = Inches(0.4)
    try:
        slide.shapes.add_picture(LOGO_PATH, left, top, width=width)
    except Exception:
        pass


def add_footer(slide, text="SLASHR  |  Agence Search Marketing  |  Avril 2026  |  Confidentiel"):
    txBox = slide.shapes.add_textbox(MARGIN_L, Inches(7.05), CONTENT_W, Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(8)
    p.font.color.rgb = WHITE_50
    p.font.name = "Inter"
    p.alignment = PP_ALIGN.LEFT


def add_page_number(slide, num):
    txBox = slide.shapes.add_textbox(Inches(12.5), Inches(7.05), Inches(0.5), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(8)
    p.font.color.rgb = WHITE_50
    p.font.name = "Inter"
    p.alignment = PP_ALIGN.RIGHT


# ============================================================
# SLIDE 1 — COUVERTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide)

# Gradient blobs (simulated with colored shapes)
for x, y, w, h, color in [
    (Inches(-1), Inches(5), Inches(5), Inches(4), ORANGE),
    (Inches(3), Inches(5.5), Inches(6), Inches(3.5), MAGENTA),
    (Inches(8), Inches(4.5), Inches(5), Inches(4.5), VIOLET),
]:
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Set transparency via XML
    from pptx.oxml.ns import qn
    from lxml import etree
    spPr = shape._element.find(qn('p:spPr'))
    solidFill = spPr.find(qn('a:solidFill'))
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is not None:
        alpha = etree.SubElement(srgb, qn('a:alpha'))
        alpha.set('val', '15000')  # 15% opacity

# Logo
add_logo(slide, left=Inches(0.8), top=Inches(0.6), width=Inches(2))

# Tag
add_tag(slide, Inches(0.8), Inches(2.5), "AGENCE SEARCH MARKETING")

# Title
add_text_box(slide, Inches(0.8), Inches(3.2), Inches(10), Inches(2),
             "Prenez le controle\ndu Search.", font_size=60, bold=True,
             font_name="Sora", line_spacing=1.05)

# Subtitle
add_text_box(slide, Inches(0.8), Inches(5.3), Inches(8), Inches(0.5),
             "Presentation a Carlos Vicente  |  Ahold Delhaize  |  Avril 2026",
             font_size=16, color=WHITE_70)


# ============================================================
# SLIDE 2 — LE SEARCH A CHANGE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide)
add_footer(slide)
add_page_number(slide, 2)

add_accent_bar(slide, MARGIN_L, Inches(1.3))
add_text_box(slide, MARGIN_L, Inches(1.5), Inches(10), Inches(0.8),
             "Le Search a change.", font_size=40, bold=True, font_name="Sora")

add_text_box(slide, MARGIN_L, Inches(2.6), Inches(10), Inches(1.5),
             "Les algorithmes evoluent, l'IA redistribue les cartes, les comportements de recherche\nse fragmentent. Ce qui fonctionnait il y a deux ans ne fonctionne plus.",
             font_size=18, color=WHITE_70)

# 4 KPI cards
cards = [
    ("Google", "Toujours dominant.\nMais les AI Overviews\nchangent la donne.", ORANGE),
    ("IA generative", "ChatGPT, Claude, Gemini,\nPerplexity — vos clients\nleur posent des questions.", MAGENTA),
    ("Social Search", "TikTok, YouTube, Instagram\nsont devenus des moteurs\nde recherche.", VIOLET),
    ("Zero-click", "60% des recherches Google\nn'aboutissent plus\na un clic.", ORANGE),
]

for i, (title, desc, accent) in enumerate(cards):
    x = MARGIN_L + Inches(i * 3)
    y = Inches(4.0)
    w = Inches(2.7)
    h = Inches(2.5)
    card = add_shape_bg(slide, x, y, w, h)
    # Accent line at top
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.fill.background()
    # Title
    add_text_box(slide, x + Inches(0.25), y + Inches(0.3), w - Inches(0.5), Inches(0.4),
                 title, font_size=16, bold=True, color=accent)
    # Description
    add_text_box(slide, x + Inches(0.25), y + Inches(0.8), w - Inches(0.5), Inches(1.5),
                 desc, font_size=13, color=WHITE_70, line_spacing=1.4)


# ============================================================
# SLIDE 3 — SLASHR EN BREF
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide)
add_footer(slide)
add_page_number(slide, 3)

add_accent_bar(slide, MARGIN_L, Inches(1.3))
add_text_box(slide, MARGIN_L, Inches(1.5), Inches(10), Inches(0.8),
             "SLASHR en bref.", font_size=40, bold=True, font_name="Sora")

add_text_box(slide, MARGIN_L, Inches(2.5), Inches(10), Inches(1),
             "On rend les marques visibles la ou leurs clients cherchent.\nGoogle, ChatGPT, TikTok, YouTube. On analyse les donnees du marche pour savoir ou investir.",
             font_size=18, color=WHITE_70)

# KPIs row
kpis = [
    ("2021", "Fondation", "a Lille"),
    ("3", "Associes", "fondateurs"),
    ("~30", "Clients", "actifs"),
    ("500K EUR", "CA annuel", "en croissance"),
    ("19", "Outils", "proprietaires"),
]

for i, (num, label, sub) in enumerate(kpis):
    x = MARGIN_L + Inches(i * 2.35)
    y = Inches(3.8)
    card = add_shape_bg(slide, x, y, Inches(2.1), Inches(1.8))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.2), Inches(1.8), Inches(0.6),
                 num, font_size=32, bold=True, color=ORANGE, font_name="Sora",
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.9), Inches(1.8), Inches(0.3),
                 label, font_size=14, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(1.2), Inches(1.8), Inches(0.3),
                 sub, font_size=11, color=WHITE_50,
                 alignment=PP_ALIGN.CENTER)

# Client logos text
add_text_box(slide, MARGIN_L, Inches(6.0), CONTENT_W, Inches(0.3),
             "Ils nous font confiance", font_size=12, bold=True, color=WHITE_50,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, MARGIN_L, Inches(6.4), CONTENT_W, Inches(0.4),
             "Vestiaire Collective  |  Carter Cash  |  SKEMA  |  EDHEC  |  Agryco  |  Essix  |  Alexandre Turpault  |  A demain",
             font_size=14, color=WHITE_70, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4 — CE QU'ON EST / CE QU'ON N'EST PAS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide)
add_footer(slide)
add_page_number(slide, 4)

add_accent_bar(slide, MARGIN_L, Inches(1.3))
add_text_box(slide, MARGIN_L, Inches(1.5), Inches(10), Inches(0.8),
             "Ce qu'on est. Ce qu'on n'est pas.", font_size=40, bold=True, font_name="Sora")

# Left column - CE QU'ON EST
x_left = MARGIN_L
y_start = Inches(2.8)
card_left = add_shape_bg(slide, x_left, y_start, Inches(5.5), Inches(4.2))

add_text_box(slide, x_left + Inches(0.3), y_start + Inches(0.2), Inches(5), Inches(0.4),
             "SLASHR est", font_size=16, bold=True, color=ORANGE)

items_is = [
    "Une agence Search Marketing strategique",
    "Un partenaire business, pas un executant SEO",
    "Des explorateurs qui testent avant de recommander",
    "Un pilote strategique qui orchestre les expertises",
    "Adaptables : le perimetre suit le besoin du client",
]

for j, item in enumerate(items_is):
    add_text_box(slide, x_left + Inches(0.3), y_start + Inches(0.7 + j * 0.65), Inches(5), Inches(0.5),
                 item, font_size=14, color=WHITE_70, line_spacing=1.3)
    # Bullet accent
    bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, x_left + Inches(0.3), y_start + Inches(0.82 + j * 0.65),
                                     Pt(6), Pt(6))
    bullet.fill.solid()
    bullet.fill.fore_color.rgb = ORANGE
    bullet.line.fill.background()
    # Indent text
    add_text_box(slide, x_left + Inches(0.6), y_start + Inches(0.7 + j * 0.65), Inches(4.7), Inches(0.5),
                 item, font_size=14, color=WHITE_70)

# Right column - CE QU'ON N'EST PAS
x_right = MARGIN_L + Inches(6.1)
card_right = add_shape_bg(slide, x_right, y_start, Inches(5.5), Inches(4.2))

add_text_box(slide, x_right + Inches(0.3), y_start + Inches(0.2), Inches(5), Inches(0.4),
             "SLASHR n'est pas", font_size=16, bold=True, color=WHITE_50)

items_not = [
    "Une agence SEO qui vend des positions Google",
    "Un prestataire qui deroule des audits generiques",
    "Un vendeur de promesses deconnectees de la data",
    "Un one-size-fits-all qui force le multicanal",
    "Un fournisseur de contenu a la chaine",
]

for j, item in enumerate(items_not):
    # X mark
    add_text_box(slide, x_right + Inches(0.3), y_start + Inches(0.65 + j * 0.65), Inches(0.3), Inches(0.5),
                 "x", font_size=14, color=WHITE_50, bold=True)
    add_text_box(slide, x_right + Inches(0.6), y_start + Inches(0.7 + j * 0.65), Inches(4.7), Inches(0.5),
                 item, font_size=14, color=WHITE_50)


# ============================================================
# SLIDE 5 — NOS 4 VALEURS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide)
add_footer(slide)
add_page_number(slide, 5)

add_accent_bar(slide, MARGIN_L, Inches(1.3))
add_text_box(slide, MARGIN_L, Inches(1.5), Inches(10), Inches(0.8),
             "Nos 4 valeurs.", font_size=40, bold=True, font_name="Sora")

values = [
    ("Curiosite", "On remet en question, on creuse, on teste sur nos propres sites avant de recommander. Le Search evolue en permanence. Notre avantage, c'est de ne jamais arreter d'apprendre.", ORANGE),
    ("Audace", "On ose dire ce qu'on pense, meme quand ca derange. On propose 90 jours au lieu de 12 mois d'engagement. On prefere une conviction argumentee a un consensus mou.", MAGENTA),
    ("Impact", "Pas de theorie sans resultat. Chaque recommandation est adossee a une donnee, chaque action est mesurable. On vise le resultat, pas la participation.", VIOLET),
    ("Transparence", "On partage nos hypotheses, nos projections et nos limites. Le client voit les memes chiffres que nous. Le ROI est conditionnel et on l'assume.", ORANGE_LIGHT),
]

for i, (title, desc, accent) in enumerate(values):
    x = MARGIN_L + Inches(i * 3)
    y = Inches(2.8)
    w = Inches(2.7)
    h = Inches(4.0)
    card = add_shape_bg(slide, x, y, w, h)
    # Number
    add_text_box(slide, x + Inches(0.25), y + Inches(0.2), Inches(0.5), Inches(0.5),
                 f"0{i+1}", font_size=28, bold=True, color=accent, font_name="Sora")
    # Title
    add_text_box(slide, x + Inches(0.25), y + Inches(0.85), w - Inches(0.5), Inches(0.4),
                 title, font_size=18, bold=True, color=WHITE)
    # Description
    add_text_box(slide, x + Inches(0.25), y + Inches(1.4), w - Inches(0.5), Inches(2.3),
                 desc, font_size=12, color=WHITE_70, line_spacing=1.5)


# ============================================================
# SLIDE 6 — NOTRE METHODE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide)
add_footer(slide)
add_page_number(slide, 6)

add_accent_bar(slide, MARGIN_L, Inches(1.3))
add_text_box(slide, MARGIN_L, Inches(1.5), Inches(10), Inches(0.8),
             "Notre methode.", font_size=40, bold=True, font_name="Sora")

add_text_box(slide, MARGIN_L, Inches(2.5), Inches(10), Inches(0.5),
             "Trois temps, toujours dans cet ordre.", font_size=18, color=WHITE_70)

steps = [
    ("01", "ANALYSER", "Comprendre le marche, les intentions de recherche et les opportunites. Decrypter avant de decider.\n\nCartographie complete de votre marche Search. Pas un echantillon — l'exhaustivite.", ORANGE),
    ("02", "PRIORISER", "Selectionner les bons leviers selon le marche, la maturite de la marque et les moyens. Dire non a ce qui ne sert pas.\n\nOn ne travaille pas tout en meme temps. On concentre les ressources.", MAGENTA),
    ("03", "EXECUTER", "Deployer, mesurer, ajuster. Des recommandations actionnables, des resultats tracables, un bilan honnete.\n\nQuick wins a 90 jours. Bilan sur donnees reelles. On continue ou on arrete.", VIOLET),
]

for i, (num, title, desc, accent) in enumerate(steps):
    x = MARGIN_L + Inches(i * 3.9)
    y = Inches(3.4)
    w = Inches(3.6)
    h = Inches(3.6)
    card = add_shape_bg(slide, x, y, w, h)
    # Top accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.fill.background()
    # Number
    add_text_box(slide, x + Inches(0.25), y + Inches(0.25), Inches(1), Inches(0.5),
                 num, font_size=36, bold=True, color=accent, font_name="Sora")
    # Title
    add_text_box(slide, x + Inches(0.25), y + Inches(0.85), w - Inches(0.5), Inches(0.4),
                 title, font_size=18, bold=True, color=WHITE)
    # Desc
    add_text_box(slide, x + Inches(0.25), y + Inches(1.4), w - Inches(0.5), Inches(2),
                 desc, font_size=11, color=WHITE_70, line_spacing=1.5)

# Arrow connectors between cards (simple)
for i in range(2):
    x = MARGIN_L + Inches((i+1) * 3.9 - 0.2)
    y = Inches(5.0)
    add_text_box(slide, x, y, Inches(0.3), Inches(0.3), ">", font_size=20, color=WHITE_50,
                 bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7 — MODELE D'ENGAGEMENT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide)
add_footer(slide)
add_page_number(slide, 7)

add_accent_bar(slide, MARGIN_L, Inches(1.3))
add_text_box(slide, MARGIN_L, Inches(1.5), Inches(10), Inches(0.8),
             "Modele d'engagement.", font_size=40, bold=True, font_name="Sora")

add_text_box(slide, MARGIN_L, Inches(2.4), Inches(10), Inches(0.5),
             "Vous restez parce que ca marche, pas parce qu'un contrat vous y oblige.",
             font_size=18, color=WHITE_70)

# Phase 1
x1 = MARGIN_L
y_p = Inches(3.3)
card1 = add_shape_bg(slide, x1, y_p, Inches(5.5), Inches(3.8))
line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x1, y_p, Inches(5.5), Pt(3))
line1.fill.solid()
line1.fill.fore_color.rgb = ORANGE
line1.line.fill.background()

add_text_box(slide, x1 + Inches(0.3), y_p + Inches(0.3), Inches(5), Inches(0.4),
             "PHASE 1 — Mission structurante", font_size=20, bold=True, color=ORANGE, font_name="Sora")
add_text_box(slide, x1 + Inches(0.3), y_p + Inches(0.85), Inches(5), Inches(0.3),
             "90 jours  |  Ponctuelle", font_size=13, bold=True, color=WHITE_50)

phase1_items = [
    "Audit SEO : cartographie complete du marche Search",
    "Benchmark : analyse des concurrents directs",
    "Architecture : arborescence et maillage optimises",
    "Contenus : specification des pages piliers",
    "Quick wins : premiers gains a 90 jours",
    "Bilan : on continue, on ajuste ou on arrete",
]
for j, item in enumerate(phase1_items):
    add_text_box(slide, x1 + Inches(0.6), y_p + Inches(1.3 + j * 0.37), Inches(4.7), Inches(0.35),
                 item, font_size=12, color=WHITE_70)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x1 + Inches(0.4), y_p + Inches(1.4 + j * 0.37), Pt(5), Pt(5))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ORANGE
    dot.line.fill.background()

# Phase 2
x2 = MARGIN_L + Inches(6.1)
card2 = add_shape_bg(slide, x2, y_p, Inches(5.5), Inches(3.8))
line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x2, y_p, Inches(5.5), Pt(3))
line2.fill.solid()
line2.fill.fore_color.rgb = MAGENTA
line2.line.fill.background()

add_text_box(slide, x2 + Inches(0.3), y_p + Inches(0.3), Inches(5), Inches(0.4),
             "PHASE 2 — Orchestration mensuelle", font_size=20, bold=True, color=MAGENTA, font_name="Sora")
add_text_box(slide, x2 + Inches(0.3), y_p + Inches(0.85), Inches(5), Inches(0.3),
             "Mensuel  |  Sans engagement", font_size=13, bold=True, color=WHITE_50)

phase2_items = [
    "Pilotage : gouvernance et ajustements strategiques",
    "Production : contenu, optimisations, link building",
    "Monitoring : positions, alertes IA, veille concurrentielle",
    "Reporting : dashboard live, revue mensuelle",
    "Comite bi-mensuel avec vos equipes",
    "Bilan trimestriel sur donnees reelles",
]
for j, item in enumerate(phase2_items):
    add_text_box(slide, x2 + Inches(0.6), y_p + Inches(1.3 + j * 0.37), Inches(4.7), Inches(0.35),
                 item, font_size=12, color=WHITE_70)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x2 + Inches(0.4), y_p + Inches(1.4 + j * 0.37), Pt(5), Pt(5))
    dot.fill.solid()
    dot.fill.fore_color.rgb = MAGENTA
    dot.line.fill.background()


# ============================================================
# SLIDE 8 — CE QUI NOUS DIFFERENCIE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide)
add_footer(slide)
add_page_number(slide, 8)

add_accent_bar(slide, MARGIN_L, Inches(1.3))
add_text_box(slide, MARGIN_L, Inches(1.5), Inches(10), Inches(0.8),
             "Ce qui nous differencie.", font_size=40, bold=True, font_name="Sora")

# Two columns comparison
x_left = MARGIN_L
x_right = MARGIN_L + Inches(6.1)
y_comp = Inches(2.8)

# Left - L'agence classique
card_l = add_shape_bg(slide, x_left, y_comp, Inches(5.5), Inches(2.5), color=SURFACE_ALT)
add_text_box(slide, x_left + Inches(0.3), y_comp + Inches(0.2), Inches(5), Inches(0.4),
             "L'agence classique dit :", font_size=16, bold=True, color=WHITE_50)
add_text_box(slide, x_left + Inches(0.3), y_comp + Inches(0.8), Inches(5), Inches(1.2),
             "\"Il y a X recherches, positionnez-vous dessus.\"\n\nApproche keyword-first. Audit generique.\nRapport et disparition.",
             font_size=14, color=WHITE_50, line_spacing=1.5)

# Right - SLASHR
card_r = add_shape_bg(slide, x_right, y_comp, Inches(5.5), Inches(2.5))
line_r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_right, y_comp, Inches(5.5), Pt(3))
line_r.fill.solid()
line_r.fill.fore_color.rgb = ORANGE
line_r.line.fill.background()
add_text_box(slide, x_right + Inches(0.3), y_comp + Inches(0.2), Inches(5), Inches(0.4),
             "SLASHR dit :", font_size=16, bold=True, color=ORANGE)
add_text_box(slide, x_right + Inches(0.3), y_comp + Inches(0.8), Inches(5), Inches(1.2),
             "\"Il y a 3 types de demande sur votre marche,\nvoici comment les adresser.\"\n\nApproche data-first. Strategie sur-mesure.\nPartnership continu.",
             font_size=14, color=WHITE_70, line_spacing=1.5)

# Bottom: 4 differentiators
diffs = [
    ("Data avant le call", "On analyse vos donnees\navant de vous parler."),
    ("ROI mesurable", "Quick wins a 90 jours.\nPas de promesse vide."),
    ("IA integree", "19 outils proprietaires.\nMonitoring IA avec Janus."),
    ("Sans engagement", "Phase 1 de 90 jours.\nVous decidez de la suite."),
]

for i, (title, desc) in enumerate(diffs):
    x = MARGIN_L + Inches(i * 3)
    y = Inches(5.7)
    card = add_shape_bg(slide, x, y, Inches(2.7), Inches(1.3))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(2.3), Inches(0.3),
                 title, font_size=13, bold=True, color=ORANGE)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.55), Inches(2.3), Inches(0.7),
                 desc, font_size=11, color=WHITE_70, line_spacing=1.4)


# ============================================================
# SLIDE 9 — JANUS (outils)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide)
add_footer(slide)
add_page_number(slide, 9)

add_accent_bar(slide, MARGIN_L, Inches(1.3))
add_text_box(slide, MARGIN_L, Inches(1.5), Inches(8), Inches(0.8),
             "Janus — Visibilite IA.", font_size=40, bold=True, font_name="Sora")

add_text_box(slide, MARGIN_L, Inches(2.5), Inches(7), Inches(1.5),
             "Savez-vous ce que l'IA dit de votre marque ?\n\nChatGPT, Claude, Gemini, Perplexity — vos clients leur posent des questions.\nJanus vous montre les reponses.\n\nVisibilite globale, sentiment moyen, sources citees,\nvisibilite par provider — tout est mesure, tout est suivi.",
             font_size=16, color=WHITE_70, line_spacing=1.5)

# Add Janus screenshot
try:
    slide.shapes.add_picture(JANUS_PATH, Inches(8.5), Inches(1.5), height=Inches(5.5))
except Exception:
    pass

# KPI cards for Janus
janus_kpis = [
    ("Visibilite globale", "Score de presence de\nvotre marque dans les IA"),
    ("Sentiment moyen", "Tonalite des reponses\ngenerees par les LLMs"),
    ("Sources citees", "Domaines et pages\nreferences par les IA"),
]

for i, (title, desc) in enumerate(janus_kpis):
    x = MARGIN_L + Inches(i * 2.5)
    y = Inches(5.5)
    card = add_shape_bg(slide, x, y, Inches(2.2), Inches(1.3))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(1.9), Inches(0.3),
                 title, font_size=11, bold=True, color=ORANGE)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.55), Inches(1.9), Inches(0.6),
                 desc, font_size=10, color=WHITE_70, line_spacing=1.4)


# ============================================================
# SLIDE 10 — CAS CLIENT : E-COMMERCE ALIMENTAIRE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide)
add_footer(slide)
add_page_number(slide, 10)

add_accent_bar(slide, MARGIN_L, Inches(1.3))
add_text_box(slide, MARGIN_L, Inches(1.5), Inches(10), Inches(0.8),
             "Cas client — E-commerce alimentaire.", font_size=40, bold=True, font_name="Sora")

add_tag(slide, MARGIN_L + Inches(8), Inches(1.65), "MARQUE PATRIMONIALE  |  B2C", bg_color=SURFACE)

# Situation
add_text_box(slide, MARGIN_L, Inches(2.6), Inches(5.5), Inches(0.3),
             "SITUATION INITIALE", font_size=12, bold=True, color=WHITE_50)
add_text_box(slide, MARGIN_L, Inches(3.0), Inches(5.5), Inches(1.5),
             "92% du trafic organique = marque. Zero acquisition generique.\nLes concurrents captent 4x plus de trafic hors-marque.\nCA 30-50M EUR.",
             font_size=14, color=WHITE_70, line_spacing=1.5)

# Actions
add_text_box(slide, MARGIN_L, Inches(4.4), Inches(5.5), Inches(0.3),
             "CE QU'ON A FAIT", font_size=12, bold=True, color=WHITE_50)
add_text_box(slide, MARGIN_L, Inches(4.8), Inches(5.5), Inches(1),
             "Architecture Search en 4 clusters thematiques.\nContenu saisonnier + donnees structurees Product/Recipe.",
             font_size=14, color=WHITE_70, line_spacing=1.5)

# Results - KPI cards
results = [
    ("x3.8", "trafic\nhors-marque"),
    ("+180K EUR", "CA organique\nattribuable"),
    ("47", "keywords\ntop 10"),
    ("12 mois", "timeline\ntotale"),
]

for i, (num, label) in enumerate(results):
    x = Inches(7.5) + Inches(i % 2 * 2.8)
    y = Inches(2.6) + Inches(i // 2 * 2.0)
    card = add_shape_bg(slide, x, y, Inches(2.5), Inches(1.7))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.2), Inches(2.2), Inches(0.6),
                 num, font_size=32, bold=True, color=ORANGE, font_name="Sora",
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.9), Inches(2.2), Inches(0.6),
                 label, font_size=12, color=WHITE_70, alignment=PP_ALIGN.CENTER, line_spacing=1.3)

# Verbatim
add_text_box(slide, MARGIN_L, Inches(6.0), Inches(10), Inches(0.8),
             "\"On pensait que notre notoriete suffisait. Les donnees nous ont montre le potentiel\nqu'on laissait aux concurrents.\"  — Directrice Marketing",
             font_size=13, color=WHITE_50, font_name="Inter")


# ============================================================
# SLIDE 11 — CAS CLIENT : RETAIL MULTI-SITES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide)
add_footer(slide)
add_page_number(slide, 11)

add_accent_bar(slide, MARGIN_L, Inches(1.3))
add_text_box(slide, MARGIN_L, Inches(1.5), Inches(10), Inches(0.8),
             "Cas client — Retail multi-sites.", font_size=40, bold=True, font_name="Sora")

add_tag(slide, MARGIN_L + Inches(8), Inches(1.65), "RESEAU PHYSIQUE  |  LOCAL", bg_color=SURFACE)

# Situation
add_text_box(slide, MARGIN_L, Inches(2.6), Inches(5.5), Inches(0.3),
             "SITUATION INITIALE", font_size=12, bold=True, color=WHITE_50)
add_text_box(slide, MARGIN_L, Inches(3.0), Inches(5.5), Inches(1.5),
             "Fiches Google My Business mal gerees. Zero strategie locale.\nConcurrents captent 3x plus de visibilite sur les requetes\n\"[produit] + [ville]\". CA 80-150M EUR, 40-80 magasins.",
             font_size=14, color=WHITE_70, line_spacing=1.5)

# Actions
add_text_box(slide, MARGIN_L, Inches(4.4), Inches(5.5), Inches(0.3),
             "CE QU'ON A FAIT", font_size=12, bold=True, color=WHITE_50)
add_text_box(slide, MARGIN_L, Inches(4.8), Inches(5.5), Inches(1),
             "Audit local complet + architecture de pages locales\n+ strategie GMB + schema LocalBusiness.",
             font_size=14, color=WHITE_70, line_spacing=1.5)

# Results
results2 = [
    ("+65%", "impressions\nGoogle Maps"),
    ("+38%", "clics\n\"itineraire\""),
    ("12 villes", "en top 3\nlocal pack"),
    ("9 mois", "timeline\ntotale"),
]

for i, (num, label) in enumerate(results2):
    x = Inches(7.5) + Inches(i % 2 * 2.8)
    y = Inches(2.6) + Inches(i // 2 * 2.0)
    card = add_shape_bg(slide, x, y, Inches(2.5), Inches(1.7))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.2), Inches(2.2), Inches(0.6),
                 num, font_size=32, bold=True, color=ORANGE, font_name="Sora",
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.9), Inches(2.2), Inches(0.6),
                 label, font_size=12, color=WHITE_70, alignment=PP_ALIGN.CENTER, line_spacing=1.3)

# Verbatim
add_text_box(slide, MARGIN_L, Inches(6.0), Inches(10), Inches(0.8),
             "\"On ne savait meme pas qu'on etait invisible sur Google Maps dans la moitie\nde nos villes.\"  — Directeur Digital",
             font_size=13, color=WHITE_50)


# ============================================================
# SLIDE 12 — CAS CLIENT : E-COMMERCE INTERNATIONAL (MIGRATION)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide)
add_footer(slide)
add_page_number(slide, 12)

add_accent_bar(slide, MARGIN_L, Inches(1.3))
add_text_box(slide, MARGIN_L, Inches(1.5), Inches(10), Inches(0.8),
             "Cas client — Migration e-commerce.", font_size=40, bold=True, font_name="Sora")

add_tag(slide, MARGIN_L + Inches(8), Inches(1.65), "INTERNATIONAL  |  REFONTE", bg_color=SURFACE)

# Situation
add_text_box(slide, MARGIN_L, Inches(2.6), Inches(5.5), Inches(0.3),
             "SITUATION INITIALE", font_size=12, bold=True, color=WHITE_50)
add_text_box(slide, MARGIN_L, Inches(3.0), Inches(5.5), Inches(1.5),
             "Refonte site prevue sans strategie SEO.\nMigration risquee : 3 000 URLs a rediriger.\nTrafic organique = 35% du CA, a preserver. 5 marches.",
             font_size=14, color=WHITE_70, line_spacing=1.5)

# Actions
add_text_box(slide, MARGIN_L, Inches(4.4), Inches(5.5), Inches(0.3),
             "CE QU'ON A FAIT", font_size=12, bold=True, color=WHITE_50)
add_text_box(slide, MARGIN_L, Inches(4.8), Inches(5.5), Inches(1),
             "Audit pre-migration + plan de redirection 1:1\n+ architecture hreflang + monitoring post-migration.",
             font_size=14, color=WHITE_70, line_spacing=1.5)

# Results
results3 = [
    ("0%", "perte de\ntrafic"),
    ("+18%", "trafic\nhors-marque"),
    ("3 000", "URLs\nredirigees"),
    ("6 mois", "timeline\ntotale"),
]

for i, (num, label) in enumerate(results3):
    x = Inches(7.5) + Inches(i % 2 * 2.8)
    y = Inches(2.6) + Inches(i // 2 * 2.0)
    card = add_shape_bg(slide, x, y, Inches(2.5), Inches(1.7))
    add_text_box(slide, x + Inches(0.15), y + Inches(0.2), Inches(2.2), Inches(0.6),
                 num, font_size=32, bold=True, color=ORANGE, font_name="Sora",
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.9), Inches(2.2), Inches(0.6),
                 label, font_size=12, color=WHITE_70, alignment=PP_ALIGN.CENTER, line_spacing=1.3)

# Verbatim
add_text_box(slide, MARGIN_L, Inches(6.0), Inches(10), Inches(0.8),
             "\"On avait peur de perdre notre trafic avec la refonte. Non seulement on l'a garde,\nmais on a gagne du terrain sur le generique.\"  — Head of E-commerce",
             font_size=13, color=WHITE_50)

# Note about Vestiaire Collective & Carter Cash
add_text_box(slide, MARGIN_L, Inches(6.7), Inches(10), Inches(0.3),
             "Expertise migration confirmee par les projets Vestiaire Collective et Carter Cash.",
             font_size=11, bold=True, color=ORANGE)


# ============================================================
# SLIDE 13 — L'EQUIPE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide)
add_footer(slide)
add_page_number(slide, 13)

add_accent_bar(slide, MARGIN_L, Inches(1.3))
add_text_box(slide, MARGIN_L, Inches(1.5), Inches(10), Inches(0.8),
             "L'equipe.", font_size=40, bold=True, font_name="Sora")

add_text_box(slide, MARGIN_L, Inches(2.3), Inches(10), Inches(0.5),
             "3 associes-fondateurs. Une equipe de consultants specialises. Un interlocuteur dedie.",
             font_size=16, color=WHITE_70)

# Associes row
associes = [
    ("Quentin Clement", "Directeur commercial\n& Associe"),
    ("Benoit Demonchaux", "Directeur de production\n& Associe"),
    ("Anthony Lecas", "Directeur Conseil\n& Associe"),
]

for i, (name, role) in enumerate(associes):
    x = MARGIN_L + Inches(i * 4)
    y = Inches(3.2)
    card = add_shape_bg(slide, x, y, Inches(3.6), Inches(1.5))
    # Avatar placeholder
    avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.25), Inches(0.8), Inches(0.8))
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = SURFACE_ALT
    avatar.line.fill.background()
    initials = "".join([w[0] for w in name.split()])
    add_text_box(slide, x + Inches(0.2), y + Inches(0.4), Inches(0.8), Inches(0.5),
                 initials, font_size=16, bold=True, color=ORANGE, alignment=PP_ALIGN.CENTER)
    # Name & role
    add_text_box(slide, x + Inches(1.2), y + Inches(0.2), Inches(2.2), Inches(0.4),
                 name, font_size=15, bold=True, color=WHITE)
    add_text_box(slide, x + Inches(1.2), y + Inches(0.6), Inches(2.2), Inches(0.6),
                 role, font_size=12, color=WHITE_70, line_spacing=1.3)

# Team members
team = [
    ("Lucas Colin", "Chef de projet senior"),
    ("Jessica Tan", "Chef de projet senior"),
    ("Pierre-Antoine Henneaux", "Consultant SEO/GEO"),
    ("Tom Chemin", "Consultant SEO/GEO"),
    ("Maxime Legru", "Consultant SEO/GEO"),
    ("Hubert Pajot", "Consultant SEA/SMA"),
]

for i, (name, role) in enumerate(team):
    col = i % 3
    row = i // 3
    x = MARGIN_L + Inches(col * 4)
    y = Inches(5.1) + Inches(row * 1.1)
    card = add_shape_bg(slide, x, y, Inches(3.6), Inches(0.9))
    # Small avatar
    avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.15), Inches(0.5), Inches(0.5))
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = SURFACE_ALT
    avatar.line.fill.background()
    initials = "".join([w[0] for w in name.split()])
    add_text_box(slide, x + Inches(0.15), y + Inches(0.22), Inches(0.5), Inches(0.4),
                 initials, font_size=11, bold=True, color=VIOLET, alignment=PP_ALIGN.CENTER)
    # Name & role
    add_text_box(slide, x + Inches(0.8), y + Inches(0.12), Inches(2.6), Inches(0.3),
                 name, font_size=13, bold=True, color=WHITE)
    add_text_box(slide, x + Inches(0.8), y + Inches(0.45), Inches(2.6), Inches(0.3),
                 role, font_size=11, color=WHITE_50)


# ============================================================
# SLIDE 14 — NOS REFERENCES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_logo(slide)
add_footer(slide)
add_page_number(slide, 14)

add_accent_bar(slide, MARGIN_L, Inches(1.3))
add_text_box(slide, MARGIN_L, Inches(1.5), Inches(10), Inches(0.8),
             "Nos references.", font_size=40, bold=True, font_name="Sora")

add_text_box(slide, MARGIN_L, Inches(2.4), Inches(10), Inches(0.5),
             "Des marques de tous secteurs qui nous font confiance pour leur strategie Search.",
             font_size=16, color=WHITE_70)

# Client grid
clients = [
    ("Vestiaire Collective", "E-commerce mode\nMigration internationale"),
    ("Carter Cash", "E-commerce auto\nMigration technique"),
    ("SKEMA", "Business School\nStrategie SEO"),
    ("EDHEC", "Business School\nVisibilite Search"),
    ("Agryco", "Agroalimentaire\nAcquisition organique"),
    ("Essix", "Orthodontie\nStrategie Search"),
    ("Alexandre Turpault", "Luxe / Linge de maison\nE-commerce SEO"),
    ("A demain", "E-commerce responsable\nStrategie Search"),
]

for i, (name, desc) in enumerate(clients):
    col = i % 4
    row = i // 4
    x = MARGIN_L + Inches(col * 3)
    y = Inches(3.3) + Inches(row * 2.2)
    card = add_shape_bg(slide, x, y, Inches(2.7), Inches(1.8))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.3), Inches(0.4),
                 name, font_size=15, bold=True, color=WHITE)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.8), Inches(2.3), Inches(0.8),
                 desc, font_size=11, color=WHITE_50, line_spacing=1.4)


# ============================================================
# SLIDE 15 — SLIDE DE CLOTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

# Gradient blobs (same as cover)
for x_pos, y_pos, w, h, color in [
    (Inches(-1), Inches(5), Inches(5), Inches(4), ORANGE),
    (Inches(3), Inches(5.5), Inches(6), Inches(3.5), MAGENTA),
    (Inches(8), Inches(4.5), Inches(5), Inches(4.5), VIOLET),
]:
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x_pos, y_pos, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    from pptx.oxml.ns import qn
    from lxml import etree
    spPr = shape._element.find(qn('p:spPr'))
    solidFill = spPr.find(qn('a:solidFill'))
    srgb = solidFill.find(qn('a:srgbClr'))
    if srgb is not None:
        alpha = etree.SubElement(srgb, qn('a:alpha'))
        alpha.set('val', '15000')

# Logo centered
add_logo(slide, left=Inches(5.4), top=Inches(1.5), width=Inches(2.5))

# Tagline
add_text_box(slide, Inches(2), Inches(3.2), Inches(9.3), Inches(1),
             "Prenez le controle du Search.", font_size=48, bold=True,
             font_name="Sora", alignment=PP_ALIGN.CENTER)

# Contact info
add_text_box(slide, Inches(2), Inches(4.5), Inches(9.3), Inches(0.5),
             "Agence Search Marketing  |  Lille, France", font_size=16, color=WHITE_70,
             alignment=PP_ALIGN.CENTER)

# Separator
add_gradient_bar(slide, Inches(5.2), Inches(5.2), Inches(3), Pt(3))

# CTA
add_text_box(slide, Inches(2), Inches(5.5), Inches(9.3), Inches(0.5),
             "Discutons de votre strategie Search.", font_size=20, color=WHITE,
             alignment=PP_ALIGN.CENTER, bold=True)

add_text_box(slide, Inches(2), Inches(6.2), Inches(9.3), Inches(0.5),
             "Quentin Clement  |  quentin@agence-slashr.fr", font_size=14, color=WHITE_50,
             alignment=PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
output_path = "/Users/quentin/Desktop/SLASHR-Presentation-Ahold-Delhaize.pptx"
prs.save(output_path)
print(f"Presentation saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
