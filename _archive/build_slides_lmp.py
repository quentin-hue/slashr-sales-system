#!/usr/bin/env python3
"""Build SLASHR SEO/SEA slides for La Mère Poulard AO response.

5 core slides (DG-level) + 8 annexe slides (data detail).
16:9, white background, collective deck format.
Data sourced from SDB.md (deal 560).
"""

import io, math
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Dimensions ──
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

# ── Colors ──
BG       = RGBColor(0xFF, 0xFF, 0xFF)
DARK     = RGBColor(0x1A, 0x1A, 0x1A)
CARD     = RGBColor(0x1A, 0x1A, 0x1A)
CARD_LT  = RGBColor(0xF3, 0xF4, 0xF6)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
RED      = RGBColor(0xE7, 0x46, 0x01)
MAGENTA  = RGBColor(0xCE, 0x08, 0xA9)
VIOLET   = RGBColor(0x89, 0x62, 0xFD)
GREEN    = RGBColor(0x22, 0xC5, 0x5E)
GRAY2    = RGBColor(0x55, 0x55, 0x55)
GRAY3    = RGBColor(0x99, 0x99, 0x99)
LTGRAY   = RGBColor(0xB3, 0xB3, 0xB3)
BAR_BG   = RGBColor(0xE5, 0xE7, 0xEB)

FONT = "Inter"


# ── Helpers ──

def set_bg(sl):
    f = sl.background.fill; f.solid(); f.fore_color.rgb = BG

def _tb(sl, l, t, w, h, txt, sz=14, c=DARK, b=False, al=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt
    p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b
    p.font.name = FONT; p.alignment = al
    return tb

def _card(sl, l, t, w, h, fill=CARD):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background(); s.adjustments[0] = 0.06
    return s

def _bar(sl, l, t, w, h, c=RED):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c
    s.line.fill.background(); s.adjustments[0] = 0.5
    return s

def _rect(sl, l, t, w, h, c=RED):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()
    return s

def _hline(sl, l, t, w, c=RED):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Emu(12700))
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background()

def _kpi(sl, l, t, w, h, val, label, accent=RED):
    _card(sl, l, t, w, h, CARD)
    _tb(sl, l+Inches(.25), t+Inches(.15), w-Inches(.5), Inches(.5),
        val, sz=26, c=accent, b=True, al=PP_ALIGN.CENTER)
    _tb(sl, l+Inches(.25), t+h-Inches(.55), w-Inches(.5), Inches(.45),
        label, sz=10, c=GRAY3, al=PP_ALIGN.CENTER)

def _title(sl, txt, y=.5, sz=36):
    _tb(sl, Inches(.8), Inches(y), Inches(11.5), Inches(1.2), txt, sz=sz, c=DARK, b=True)

def _footer(sl, txt):
    _tb(sl, Inches(.8), Inches(6.8), Inches(11.5), Inches(.2), txt, sz=8, c=GRAY3)

def _bullet(sl, x, y, w, txt, sz=14, c=DARK):
    _tb(sl, x, y, Inches(.25), Inches(.3), "●", sz=8, c=RED, al=PP_ALIGN.LEFT)
    _tb(sl, x+Inches(.3), y-Inches(.02), w-Inches(.3), Inches(.5), txt, sz=sz, c=c)

def _annexe_label(sl, num, title):
    _tb(sl, Inches(.8), Inches(.12), Inches(6), Inches(.25),
        f"ANNEXE A{num} · {title}", sz=9, c=GRAY3, b=True)


# ═══════════════════════════════════════════════════
#  CORE SLIDE 1 — LE SEARCH CHANGE
# ═══════════════════════════════════════════════════

def slide_geo_core(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    _title(sl, "Le Search change.", sz=44)
    _tb(sl, Inches(.8), Inches(1.4), Inches(11.5), Inches(.5),
        "Les IA répondent à la place de Google.", sz=22, c=GRAY2)

    # 3 platform KPIs — simple, macro
    cw, ch, cy = Inches(3.5), Inches(1.6), Inches(2.4)
    gap = Inches(.37)

    platforms = [
        ("ChatGPT",             "740 M",  "utilisateurs/mois",    VIOLET),
        ("Google AI Overviews", "1 Md+",  "recherches avec IA/jour", RED),
        ("Perplexity",          "100 M+", "requêtes/mois",        MAGENTA),
    ]
    for i, (name, val, unit, accent) in enumerate(platforms):
        x = Inches(.8) + i * (cw + gap)
        _card(sl, x, cy, cw, ch, CARD)
        _tb(sl, x+Inches(.2), cy+Inches(.15), cw-Inches(.4), Inches(.25),
            name, sz=11, c=WHITE, b=True)
        _tb(sl, x+Inches(.2), cy+Inches(.45), cw-Inches(.4), Inches(.5),
            val, sz=32, c=accent, b=True)
        _tb(sl, x+Inches(.2), cy+Inches(1.05), cw-Inches(.4), Inches(.25),
            unit, sz=10, c=GRAY3)

    # Key message
    _tb(sl, Inches(.8), Inches(4.5), Inches(11.5), Inches(.6),
        'Quand un internaute demande "meilleur coffret biscuit cadeau",\n'
        "l'IA cite une marque. Si votre contenu n'est pas structuré, ce ne sera pas la vôtre.",
        sz=16, c=DARK)

    # Advantage callout
    _card(sl, Inches(.8), Inches(5.5), Inches(11.5), Inches(.75), CARD)
    _tb(sl, Inches(1.1), Inches(5.58), Inches(2), Inches(.2),
        "VOTRE AVANTAGE", sz=10, c=GREEN, b=True)
    _tb(sl, Inches(1.1), Inches(5.82), Inches(10.5), Inches(.35),
        "1888, Mont-Saint-Michel, savoir-faire artisanal — les IA valorisent "
        "l'authenticité et l'expertise (E-E-A-T). Votre histoire est un actif naturel.",
        sz=12, c=WHITE)

    _footer(sl, "GEO = Generative Engine Optimization · Le SEO des réponses IA")


# ═══════════════════════════════════════════════════
#  CORE SLIDE 2 — LE CONSTAT
# ═══════════════════════════════════════════════════

def slide_constat(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)

    # Hero number
    _tb(sl, Inches(.8), Inches(.6), Inches(11.5), Inches(1.4),
        "+308 000", sz=80, c=DARK, b=True)
    _tb(sl, Inches(.8), Inches(2.1), Inches(11.5), Inches(.5),
        "recherches par mois sur vos produits.", sz=24, c=GRAY2)
    _tb(sl, Inches(.8), Inches(2.7), Inches(11.5), Inches(.5),
        "Votre site est invisible pour ceux qui ne vous connaissent pas.", sz=24, c=RED, b=True)

    # GSC reality callout
    _card(sl, Inches(.8), Inches(3.5), Inches(11.5), Inches(.7), CARD)
    _tb(sl, Inches(1.1), Inches(3.55), Inches(2), Inches(.2),
        "DONNÉES RÉELLES (GSC)", sz=9, c=RED, b=True)
    _tb(sl, Inches(1.1), Inches(3.78), Inches(10.5), Inches(.35),
        "<100 clics hors-marque par mois. Sur 4 500 visiteurs organiques, "
        "95 % vous connaissent déjà. Le reste du marché va ailleurs.",
        sz=13, c=WHITE)

    # Gap bar
    bx, by, bw = Inches(.8), Inches(4.5), Inches(11.5)
    _bar(sl, bx, by, bw, Inches(.35), BAR_BG)
    _bar(sl, bx, by, Inches(.23), Inches(.35), RED)  # 2% of 11.5"
    _tb(sl, bx, by+Inches(.42), Inches(2), Inches(.25),
        "2 % captés", sz=11, c=RED, b=True)
    _tb(sl, Inches(9.3), by+Inches(.42), Inches(3), Inches(.25),
        "98 % non captés", sz=11, c=GRAY3, al=PP_ALIGN.RIGHT)

    # 3 KPI cards
    cw, ch, cy = Inches(3.5), Inches(1.0), Inches(5.3)
    gap = Inches(.37)
    _kpi(sl, Inches(.8), cy, cw, ch,
         "186 000", "Recherches commerciales\nbiscuits, coffrets, cadeaux", RED)
    _kpi(sl, Inches(.8)+cw+gap, cy, cw, ch,
         "99 000", "Recherches informationnelles\nrecettes, guides, terroir", RED)
    _kpi(sl, Inches(.8)+2*(cw+gap), cy, cw, ch,
         "23 000", "Recherches internationales\n35 % de votre trafic actuel", RED)

    _footer(sl, "Sources : DataForSEO (marché) · GSC 12 mois (trafic réel) · 2025")


# ═══════════════════════════════════════════════════
#  CORE SLIDE 2 — POURQUOI
# ═══════════════════════════════════════════════════

def slide_pourquoi(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    _title(sl, "Google vous fait confiance.\nMais votre site ne lui donne rien à montrer.", sz=34)

    # ── Left: L'ATOUT ──
    ax, ay, aw, ah = Inches(.8), Inches(2.0), Inches(5.2), Inches(2.5)
    _card(sl, ax, ay, aw, ah, CARD)
    _rect(sl, ax, ay, Inches(.12), ah, GREEN)  # green left border
    _tb(sl, ax+Inches(.35), ay+Inches(.15), Inches(2), Inches(.25),
        "L'ATOUT", sz=10, c=GREEN, b=True)

    _tb(sl, ax+Inches(.35), ay+Inches(.45), aw-Inches(.7), Inches(.5),
        '#3 sur "biscuits"', sz=32, c=WHITE, b=True)
    _tb(sl, ax+Inches(.35), ay+Inches(.95), aw-Inches(.7), Inches(.25),
        "40 500 recherches/mois", sz=16, c=RED, b=True)

    _hline(sl, ax+Inches(.35), ay+Inches(1.3), aw-Inches(.7), RGBColor(0x33,0x33,0x33))

    _tb(sl, ax+Inches(.35), ay+Inches(1.4), aw-Inches(.7), Inches(.5),
        "+27 % en 2025", sz=32, c=GREEN, b=True)
    _tb(sl, ax+Inches(.35), ay+Inches(1.85), aw-Inches(.7), Inches(.5),
        "55K€ de CA web sans stratégie digitale.\n"
        "La marque porte seule — mais ce plafond approche.",
        sz=11, c=LTGRAY)

    # ── Right: 3 VERROUS ──
    vx, vw = Inches(6.3), Inches(6)
    vh = Inches(.7)
    vgap = Inches(.15)

    _tb(sl, vx, Inches(1.85), Inches(3), Inches(.25),
        "3 VERROUS", sz=10, c=RED, b=True)

    # Verrou 1
    vy1 = Inches(2.2)
    _card(sl, vx, vy1, vw, vh, CARD_LT)
    _tb(sl, vx+Inches(.2), vy1+Inches(.08), Inches(1.4), Inches(.3),
        "157 mots/page", sz=18, c=RED, b=True)
    _tb(sl, vx+Inches(1.8), vy1+Inches(.12), Inches(4), Inches(.25),
        "→ Google n'a rien à indexer. Le catalogue est muet.", sz=12, c=GRAY2)

    # Verrou 2
    vy2 = vy1 + vh + vgap
    _card(sl, vx, vy2, vw, vh, CARD_LT)
    _tb(sl, vx+Inches(.2), vy2+Inches(.08), Inches(1.4), Inches(.3),
        "0 contenu ciblé", sz=18, c=RED, b=True)
    _tb(sl, vx+Inches(1.8), vy2+Inches(.12), Inches(4), Inches(.25),
        "→ Recettes, coffrets, cadeaux : ces requêtes vont ailleurs.", sz=12, c=GRAY2)

    # Verrou 3
    vy3 = vy2 + vh + vgap
    _card(sl, vx, vy3, vw, vh, CARD_LT)
    _tb(sl, vx+Inches(.2), vy3+Inches(.08), Inches(1.4), Inches(.3),
        "0 € en Ads", sz=18, c=RED, b=True)
    _tb(sl, vx+Inches(1.8), vy3+Inches(.12), Inches(4), Inches(.25),
        "→ Sur un marché où le clic coûte moins d'1 €.", sz=12, c=GRAY2)

    # Bottom punchline
    _tb(sl, Inches(.8), Inches(5.2), Inches(11.5), Inches(.5),
        "L'autorité est là. La croissance est là. Le contenu manque.",
        sz=20, c=DARK, b=True)
    _tb(sl, Inches(.8), Inches(5.8), Inches(11.5), Inches(.4),
        "Vous croissez par la force de la marque. Mais <100 clics hors-marque/mois "
        "= le marché ne vous voit pas. La refonte est la fenêtre.",
        sz=14, c=GRAY2)

    _footer(sl, "Sources : DataForSEO + crawl + GSC + WooCommerce 2024-2025 · Détails en annexe A1-A6")


# ═══════════════════════════════════════════════════
#  CORE SLIDE 4 — LA CONCURRENCE
# ═══════════════════════════════════════════════════

def slide_concurrence(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    _title(sl, "Pendant ce temps, vos concurrents\ncaptent le marché.", sz=34)

    # ── Left: TRAFIC ORGANIQUE ──
    lx, lw = Inches(.8), Inches(5.4)
    _tb(sl, lx, Inches(1.85), lw, Inches(.25),
        "TRAFIC ORGANIQUE (GRATUIT)", sz=10, c=RED, b=True)
    _hline(sl, lx, Inches(2.15), Inches(3), RED)

    _tb(sl, lx, Inches(2.4), lw, Inches(.6),
        "La Trinitaine capte ×100\nplus de nouveaux clients.", sz=24, c=DARK, b=True)

    # Comparison bars
    by = Inches(3.3)
    bw_max = Inches(4.5)
    # La Trinitaine: ~16,000 non-brand (82% of 20,016 — DataForSEO)
    _tb(sl, lx, by, Inches(2.2), Inches(.25), "La Trinitaine", sz=12, c=GRAY2)
    _bar(sl, lx+Inches(2.3), by+Inches(.02), bw_max, Inches(.22), RED)
    _tb(sl, lx+Inches(2.3)+bw_max+Inches(.1), by, Inches(1.5), Inches(.25),
        "~16 000", sz=12, c=RED, b=True)

    # LMP: <100 non-brand (GSC réel)
    by2 = by + Inches(.4)
    lmp_w = bw_max * (100 / 16000)  # quasi invisible
    _tb(sl, lx, by2, Inches(2.2), Inches(.25), "La Mère Poulard", sz=12, c=DARK, b=True)
    _bar(sl, lx+Inches(2.3), by2+Inches(.02), lmp_w, Inches(.22), DARK)
    _tb(sl, lx+Inches(2.3)+lmp_w+Inches(.1), by2, Inches(1.5), Inches(.25),
        "<100", sz=12, c=DARK, b=True)

    _tb(sl, lx, by2+Inches(.4), lw, Inches(.25),
        "clics hors-marque / mois (LMP = données GSC réelles)", sz=10, c=GRAY3)

    _tb(sl, lx, Inches(4.5), lw, Inches(.5),
        "+27 % de croissance en 2025 sans aucune stratégie.\n"
        "La marque porte seule. Le contenu structuré manque.",
        sz=13, c=GRAY2)

    # ── Right: INVESTISSEMENT PUBLICITAIRE ──
    rx, rw = Inches(6.55), Inches(5.8)
    _tb(sl, rx, Inches(1.85), rw, Inches(.25),
        "INVESTISSEMENT PUBLICITAIRE", sz=10, c=VIOLET, b=True)
    _hline(sl, rx, Inches(2.15), Inches(3), VIOLET)

    competitors_paid = [
        ("L'Atelier St Michel",   "30 – 150 K€/mois", "Marque nationale,\neffet retail",      MAGENTA),
        ("La Trinitaine",         "10 – 50 K€/mois",  "Coffrets, patrimoine,\nQ4 agressif",    RED),
        ("Biscuiterie Pointe du Raz", "3 – 15 K€/mois", "Terroir Bretagne,\ncoffrets cadeaux", GRAY3),
        ("La Mère Poulard",       "0 €/mois",         "",                                       DARK),
    ]
    cy2 = Inches(2.4)
    for name, budget, note, accent in competitors_paid:
        is_lmp = "Poulard" in name
        _card(sl, rx, cy2, rw, Inches(.55), CARD if is_lmp else CARD_LT)
        tc = WHITE if is_lmp else DARK
        _tb(sl, rx+Inches(.2), cy2+Inches(.05), Inches(2.5), Inches(.25),
            name, sz=12, c=tc, b=is_lmp)
        _tb(sl, rx+Inches(2.8), cy2+Inches(.05), Inches(2.5), Inches(.25),
            budget, sz=14, c=RED if is_lmp else accent, b=True)
        cy2 += Inches(.62)

    _tb(sl, rx, Inches(4.95), rw, Inches(.35),
        "Source : estimations Cocoa (benchmark sectoriel)", sz=9, c=GRAY3)

    # Bottom
    _tb(sl, Inches(.8), Inches(5.6), Inches(11.5), Inches(.5),
        "Le marché n'attend pas. Chaque mois sans activation est un mois "
        "offert à la concurrence.",
        sz=16, c=DARK, b=True)

    _footer(sl, "Sources : GSC + WooCommerce (LMP) · DataForSEO (La Trinitaine) · Cocoa (paid estimations) · 2025")


# ═══════════════════════════════════════════════════
#  CORE SLIDE 5 — IDENTITÉ GOOGLE : 3 PILIERS
# ═══════════════════════════════════════════════════

def slide_identite(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    _title(sl, "Ce qui est hors de votre site compte\nautant que ce qui est dessus.", sz=36)

    _tb(sl, Inches(.8), Inches(1.55), Inches(11.5), Inches(.35),
        "Google évalue votre marque dans son ensemble : site, avis, réseaux, mentions, données structurées. "
        "Aujourd'hui, vos signaux off-site brouillent votre identité.",
        sz=13, c=GRAY2)

    # ── Left: 5 signal cards (current state) ──
    lx, lw = Inches(.8), Inches(5.2)
    _tb(sl, lx, Inches(2.2), lw, Inches(.25),
        "VOS SIGNAUX OFF-SITE AUJOURD'HUI", sz=10, c=RED, b=True)
    _hline(sl, lx, Inches(2.5), Inches(3), RED)

    signals = [
        ("Google Business",   "1,8★",  "25 avis",    RED,    "Restaurant = fiche visible. Biscuiterie 4,4★ noyée."),
        ("TripAdvisor",       "2,7★",  "5 578 avis", RED,    "Note restaurant pèse sur l'image globale du groupe."),
        ("Google Shopping",   "0",     "produit",    MAGENTA, "Vos coffrets existent, mais c'est La Grande Epicerie qui vend."),
        ("Réseaux sociaux",   "<1%",   "du trafic",  GRAY3,   "215 sessions/an via social. Aucun levier d'acquisition."),
        ("IA / Knowledge Graph", "?",  "",           GRAY3,   "Google confond restaurant, biscuiterie et hôtel."),
    ]

    sy = Inches(2.7)
    for name, score, unit, accent, desc in signals:
        _card(sl, lx, sy, lw, Inches(.55), CARD_LT)
        _tb(sl, lx+Inches(.15), sy+Inches(.05), Inches(1.8), Inches(.22),
            name, sz=11, c=DARK, b=True)
        _tb(sl, lx+Inches(2.0), sy+Inches(.05), Inches(.6), Inches(.22),
            score, sz=14, c=accent, b=True)
        _tb(sl, lx+Inches(2.6), sy+Inches(.08), Inches(.8), Inches(.2),
            unit, sz=9, c=GRAY3)
        _tb(sl, lx+Inches(3.5), sy+Inches(.08), Inches(1.6), Inches(.45),
            desc, sz=9, c=GRAY2)
        sy += Inches(.62)

    # ── Right: what this means ──
    rx, rw = Inches(6.3), Inches(5.8)
    _tb(sl, rx, Inches(2.2), rw, Inches(.25),
        "POURQUOI C'EST CRITIQUE", sz=10, c=VIOLET, b=True)
    _hline(sl, rx, Inches(2.5), Inches(3), VIOLET)

    # Card 1: E-E-A-T
    _card(sl, rx, Inches(2.7), rw, Inches(1.3), CARD)
    _tb(sl, rx+Inches(.25), Inches(2.8), Inches(1.2), Inches(.25),
        "E-E-A-T", sz=18, c=RED, b=True)
    _tb(sl, rx+Inches(1.6), Inches(2.82), rw-Inches(1.9), Inches(.2),
        "le critère n°1 de Google", sz=10, c=GRAY3)
    _tb(sl, rx+Inches(.25), Inches(3.15), rw-Inches(.5), Inches(.7),
        "Expérience · Expertise · Autorité · Fiabilité.\n"
        "Google ne regarde pas que vos pages. Il croise\n"
        "avis, mentions, profils sociaux, données structurées.\n"
        "Un signal contradictoire = confiance dégradée.",
        sz=11, c=GRAY3)

    # Card 2: Impact concret
    _card(sl, rx, Inches(4.15), rw, Inches(1.35), CARD)
    _tb(sl, rx+Inches(.25), Inches(4.25), rw-Inches(.5), Inches(.22),
        "Impact concret sur La Mère Poulard", sz=11, c=WHITE, b=True)
    impacts = [
        ("Classement",   "Les avis négatifs (restaurant) freinent\nle positionnement de la biscuiterie"),
        ("Conversion",   "Un internaute qui voit 1,8★ avant d'acheter\nun coffret à 50 € hésite"),
        ("IA générative", "ChatGPT/Gemini confondent vos activités\n→ réponses incorrectes sur la marque"),
    ]
    iy = Inches(4.55)
    for label, desc in impacts:
        _tb(sl, rx+Inches(.25), iy, Inches(1.2), Inches(.2),
            label, sz=10, c=RED, b=True)
        _tb(sl, rx+Inches(1.5), iy-Inches(.02), rw-Inches(1.8), Inches(.42),
            desc, sz=9, c=GRAY3)
        iy += Inches(.42)

    # ── Bottom: the 3-pillar response ──
    _card(sl, Inches(.8), Inches(5.85), Inches(11.5), Inches(.72), CARD)
    bx = Inches(1.1)
    responses = [
        ("REFONTE", RED,    "Séparer les entités (Schema, GBP, architecture) → identité claire"),
        ("SEA",     VIOLET, "Contrôler le message dès J1 (Shopping, Search) → contourner la confusion"),
        ("IA",      MAGENTA, "Structurer les données → la bonne marque dans les réponses de demain"),
    ]
    for i, (label, accent, desc) in enumerate(responses):
        x = bx + i * Inches(3.7)
        _tb(sl, x, Inches(5.92), Inches(.6), Inches(.2),
            label, sz=10, c=accent, b=True)
        _tb(sl, x+Inches(.7), Inches(5.92), Inches(2.8), Inches(.5),
            desc, sz=10, c=GRAY3)

    _footer(sl, "Sources : Google Business Profile · TripAdvisor · GA4 · mars 2026")


# ═══════════════════════════════════════════════════
#  CORE SLIDE 6 — L'APPROCHE SEO
# ═══════════════════════════════════════════════════

def slide_seo(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    _title(sl, "Profiter de la refonte pour construire\nle site que Google attend.", sz=34)

    _tb(sl, Inches(.8), Inches(1.7), Inches(11.5), Inches(.35),
        "La refonte e-commerce est la fenêtre. Sans SEO intégré dès la conception, "
        "risque de perdre 30 % du trafic existant.",
        sz=13, c=GRAY2)

    # ── 3 pillar cards ──
    cw, ch = Inches(3.5), Inches(3.1)
    cy = Inches(2.5)
    gap = Inches(.37)

    pillars = [
        ("1", "STRUCTURER", RED,
         "Organiser le site par\nintention d'achat",
         "Le site doit refléter comment\n"
         "les clients cherchent, pas\n"
         "comment vous rangez vos produits.\n\n"
         "Clusters : biscuits, coffrets cadeaux,\n"
         "recettes, terroir normand, international."),
        ("2", "CRÉER", RED,
         "Le contenu que personne\nne produit",
         "1888, savoir-faire artisanal,\n"
         "Mont-Saint-Michel.\n\n"
         "Recettes, guides, fiches enrichies.\n"
         "Votre légitimité (E-E-A-T) est un\n"
         "avantage naturel. 0 concurrent\n"
         "artisanal positionné."),
        ("3", "SÉCURISER", GREEN,
         "La migration\nde domaine",
         "Vous êtes #3 sur « biscuits ».\n"
         "Il faut le rester.\n\n"
         "Redirections 301 exhaustives,\n"
         "transfert d'autorité, monitoring\n"
         "indexation post-migration."),
    ]

    for i, (num, label, accent, headline, body) in enumerate(pillars):
        x = Inches(.8) + i * (cw + gap)
        _card(sl, x, cy, cw, ch, CARD)
        # Number badge
        _rect(sl, x, cy, Inches(.45), Inches(.38), accent)
        _tb(sl, x, cy+Inches(.02), Inches(.45), Inches(.35),
            num, sz=18, c=WHITE, b=True, al=PP_ALIGN.CENTER)
        _tb(sl, x+Inches(.55), cy+Inches(.06), cw-Inches(.8), Inches(.25),
            label, sz=11, c=accent, b=True)
        # Headline
        _tb(sl, x+Inches(.2), cy+Inches(.5), cw-Inches(.4), Inches(.55),
            headline, sz=15, c=WHITE, b=True)
        # Body
        _tb(sl, x+Inches(.2), cy+Inches(1.15), cw-Inches(.4), Inches(1.8),
            body, sz=10, c=LTGRAY)

    # Bottom result
    _tb(sl, Inches(.8), Inches(5.9), Inches(11.5), Inches(.45),
        "Résultat : un actif qui s'apprécie avec le temps. "
        "Trafic gratuit, croissant, composé.",
        sz=16, c=DARK, b=True)

    _footer(sl, "SLASHR · Search Engine Optimization · Détails en annexe A1-A6")


# ═══════════════════════════════════════════════════
#  CORE SLIDE 4 — L'APPROCHE SEA
# ═══════════════════════════════════════════════════

def slide_sea(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    _title(sl, "Activer le marché dès le premier jour.", sz=36)

    # ── Left: POURQUOI ──
    lx, ly, lw = Inches(.8), Inches(1.9), Inches(5.2)
    _tb(sl, lx, ly, lw, Inches(.25),
        "POURQUOI LE SEA EN COMPLÉMENT", sz=10, c=VIOLET, b=True)
    _hline(sl, lx, ly+Inches(.3), Inches(3.5), VIOLET)

    reasons = [
        ("Le SEO prend 6 à 12 mois.",
         "Les Ads génèrent du CA dès le jour 1."),
        ("69 500 recherches « bretonnes » /mois.",
         "LMP est normande — seul le SEA peut cibler\nces requêtes sans contrainte éditoriale."),
        ("Pics saisonniers : Noël & Pâques.",
         "Coffrets +404 %, paniers +311 %. Il faut être\nprésent 6 semaines avant chaque pic."),
    ]
    ry = ly + Inches(.5)
    for reason, detail in reasons:
        _tb(sl, lx, ry, lw, Inches(.3), reason, sz=14, c=DARK, b=True)
        _tb(sl, lx, ry+Inches(.3), lw, Inches(.45), detail, sz=11, c=GRAY2)
        ry += Inches(.85)

    # ── Right: 3 LEVIERS ──
    rx, rw, rh = Inches(6.5), Inches(5.8), Inches(.95)
    _tb(sl, rx, Inches(1.9), rw, Inches(.25),
        "LES 3 LEVIERS MÉDIA", sz=10, c=VIOLET, b=True)
    _hline(sl, rx, Inches(2.2), Inches(3.5), VIOLET)

    levers = [
        ("GOOGLE ADS", VIOLET,
         "Search brand + hors-marque\nShopping / PMax · International EN"),
        ("META ADS", MAGENTA,
         "Prospection familles 35-55 ans\nRetargeting visiteurs site"),
        ("CRM", GREEN,
         "Welcome · Abandon panier\nPost-achat · Relance Q4 CE/CSE"),
    ]
    ly2 = Inches(2.45)
    for label, accent, desc in levers:
        _card(sl, rx, ly2, rw, rh, CARD)
        _tb(sl, rx+Inches(.2), ly2+Inches(.1), Inches(1.8), Inches(.2),
            label, sz=10, c=accent, b=True)
        _tb(sl, rx+Inches(.2), ly2+Inches(.35), rw-Inches(.4), Inches(.5),
            desc, sz=11, c=LTGRAY)
        ly2 += rh + Inches(.12)

    # ── Bottom: saisonnalité callout ──
    by = Inches(5.65)
    _card(sl, Inches(.8), by, Inches(11.5), Inches(.7), CARD)
    _tb(sl, Inches(1.1), by+Inches(.08), Inches(2), Inches(.2),
        "SAISONNALITÉ", sz=10, c=RED, b=True)
    _tb(sl, Inches(1.1), by+Inches(.32), Inches(10.5), Inches(.3),
        "Budget ×2 en septembre–décembre (période cadeau + CE/CSE). "
        "Le SEA capte la demande saisonnière que le SEO seul ne peut pas anticiper.",
        sz=12, c=WHITE)

    _footer(sl, "Cocoa · Search Engine Advertising + Social + CRM")


# ═══════════════════════════════════════════════════
#  CORE SLIDE 5 — LA SYNERGIE
# ═══════════════════════════════════════════════════

def _compound_png():
    """Generate compound effect chart: SEA flat + SEO exponential."""
    months = list(range(1, 13))
    sea = [30, 33, 35, 37, 38, 39, 39, 40, 55, 65, 58, 42]
    seo = [2, 4, 7, 11, 16, 22, 29, 37, 45, 53, 61, 70]
    combined = [s + e for s, e in zip(sea, seo)]

    fig, ax = plt.subplots(figsize=(9.5, 3.2))
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    ax.fill_between(months, 0, sea, alpha=.15, color="#8962FD")
    ax.fill_between(months, sea, combined, alpha=.15, color="#E74601")
    ax.plot(months, sea, color="#8962FD", linewidth=2.5, label="SEA — CA immédiat")
    ax.plot(months, combined, color="#E74601", linewidth=2.5, label="SEO + SEA — Effet composé")

    ax.set_xticks(months)
    ax.set_xticklabels([f"M{m}" for m in months], fontsize=9, color="#999")
    ax.set_ylabel("Impact (index)", fontsize=10, color="#555")
    ax.set_xlim(.5, 12.5)
    ax.set_ylim(0, max(combined) * 1.15)
    ax.legend(loc="upper left", fontsize=10, frameon=False)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["bottom"].set_color("#ddd")
    ax.spines["left"].set_color("#ddd")
    ax.tick_params(colors="#999")

    # Phase labels
    for x1, x2, label in [(0.5, 3.5, "Fondation"), (3.5, 6.5, "Lancement"), (6.5, 12.5, "Accélération")]:
        ax.axvspan(x1, x2, alpha=.025, color="#999")
        ax.text((x1+x2)/2, max(combined)*1.07, label,
                ha="center", fontsize=8, color="#999")

    buf = io.BytesIO()
    plt.savefig(buf, format="png", dpi=180, bbox_inches="tight", facecolor="white")
    plt.close(); buf.seek(0)
    return buf


def slide_synergie(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    _title(sl, "SEO × SEA = effet composé.", sz=40, y=.3)

    # Chart
    chart = _compound_png()
    sl.shapes.add_picture(chart, Inches(.5), Inches(1.4), Inches(11.5), Inches(3.2))

    # ── Timeline: 3 phases ──
    ty = Inches(4.8)
    pw = Inches(3.5)
    gap = Inches(.37)

    phases = [
        ("Phase 1 · Fondation", "M1 – M3", RED,
         "Structure SEO intégrée à la refonte\nContenu prioritaire (6 territoires)\nPlan de migration validé\nSetup Google Ads + Shopping"),
        ("Phase 2 · Lancement", "M4 – M6", GREEN,
         "Migration sécurisée + monitoring\nActivation Ads (Search + Meta + CRM)\nContenu recettes + guides\nAnticipation pics Q4"),
        ("Phase 3 · Accélération", "M7 – M12", VIOLET,
         "International EN (SEO + Ads)\nFondations B2B (CE/CSE + Pro)\nOptimisation continue\nScale budgets selon ROAS"),
    ]

    for i, (title, period, accent, items) in enumerate(phases):
        x = Inches(.8) + i * (pw + gap)
        _hline(sl, x, ty, pw, accent)
        _tb(sl, x, ty+Inches(.08), Inches(1), Inches(.2),
            period, sz=9, c=accent, b=True)
        _tb(sl, x+Inches(1.1), ty+Inches(.08), pw-Inches(1.1), Inches(.2),
            title, sz=11, c=DARK, b=True)
        _tb(sl, x, ty+Inches(.35), pw, Inches(.9),
            items, sz=9, c=GRAY2)

    # Bottom
    _card(sl, Inches(.8), Inches(6.35), Inches(11.5), Inches(.42), CARD)
    _tb(sl, Inches(1.1), Inches(6.38), Inches(10.5), Inches(.35),
        "Le SEA finance le temps que le SEO s'installe. Le SEO réduit le CPC du SEA. "
        "Les données Ads alimentent la stratégie SEO.",
        sz=11, c=WHITE)

    _footer(sl, "SLASHR (SEO) + Cocoa (Ads) + Fractory (site) · Réponse collective AO La Mère Poulard")


# ═══════════════════════════════════════════════════
#  CORE SLIDE 8 — PROJECTION CA
# ═══════════════════════════════════════════════════

def slide_projection(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    _title(sl, "Où investir pour maximiser le retour.", sz=36)

    # ── Current state strip ──
    _tb(sl, Inches(.8), Inches(1.45), Inches(1.5), Inches(.25),
        "AUJOURD'HUI", sz=9, c=GRAY3, b=True)
    # 4 mini KPIs inline
    state_items = [
        ("CA web net", "55 K€/an (+27%)"),
        ("Trafic total", "~7 000 sessions/mois"),
        ("dont international", "35 % (GSC)"),
        ("Ads + CRM", "0 €"),
    ]
    sx = Inches(2.5)
    for label, val in state_items:
        _tb(sl, sx, Inches(1.4), Inches(1.2), Inches(.2), label, sz=8, c=GRAY3)
        _tb(sl, sx, Inches(1.58), Inches(1.2), Inches(.2), val, sz=11, c=DARK, b=True)
        sx += Inches(2.2)

    _hline(sl, Inches(.8), Inches(1.9), Inches(11.5), BAR_BG)

    # ── 3 Lever cards ──
    cw = Inches(3.5)
    gap = Inches(.37)
    cy = Inches(2.15)

    # Heights: card 1 taller (priority), others standard
    ch1, ch2, ch3 = Inches(4.15), Inches(4.15), Inches(4.15)

    # ── LEVIER 1 : SEO (BUILD + RUN) ──
    x1 = Inches(.8)
    _card(sl, x1, cy, cw, ch1, CARD)
    # Priority badge
    _rect(sl, x1, cy, cw, Inches(.32), RED)
    _tb(sl, x1, cy+Inches(.04), cw, Inches(.25),
        "PRIORITÉ 2026", sz=10, c=WHITE, b=True, al=PP_ALIGN.CENTER)

    _tb(sl, x1+Inches(.2), cy+Inches(.45), cw-Inches(.4), Inches(.25),
        "SEO", sz=16, c=WHITE, b=True)

    # BUILD phase
    _tb(sl, x1+Inches(.2), cy+Inches(.78), cw-Inches(.4), Inches(.2),
        "BUILD · dans la refonte", sz=9, c=RED, b=True)
    _tb(sl, x1+Inches(.2), cy+Inches(.98), cw-Inches(.4), Inches(.4),
        "Structure, architecture, Schema,\nmigration sécurisée", sz=10, c=LTGRAY)
    _tb(sl, x1+Inches(.2), cy+Inches(1.35), cw-Inches(.4), Inches(.2),
        "Surcoût ~0 € (intégré à la refonte)", sz=10, c=GREEN, b=True)

    # Separator
    _hline(sl, x1+Inches(.2), cy+Inches(1.65), cw-Inches(.4), RGBColor(0x33,0x33,0x33))

    # RUN phase
    _tb(sl, x1+Inches(.2), cy+Inches(1.78), cw-Inches(.4), Inches(.2),
        "RUN · accompagnement mensuel", sz=9, c=RED, b=True)
    _tb(sl, x1+Inches(.2), cy+Inches(1.98), cw-Inches(.4), Inches(.4),
        "Contenu expert, optimisation continue,\nsuivi positions, maillage", sz=10, c=LTGRAY)
    _tb(sl, x1+Inches(.2), cy+Inches(2.35), cw-Inches(.4), Inches(.2),
        "→ C'est le Run qui génère la croissance", sz=10, c=WHITE, b=True)

    # Result
    _hline(sl, x1+Inches(.2), cy+Inches(2.65), cw-Inches(.4), RED)
    _tb(sl, x1+Inches(.2), cy+Inches(2.75), Inches(2.5), Inches(.3),
        "+35 K€/an", sz=20, c=RED, b=True)
    _tb(sl, x1+Inches(.2), cy+Inches(3.05), cw-Inches(.4), Inches(.2),
        "trafic gratuit, croissant, composé", sz=10, c=LTGRAY)

    _tb(sl, x1+Inches(.2), cy+Inches(3.35), cw-Inches(.4), Inches(.55),
        "Le Build pose les fondations.\n"
        "Le Run construit l'actif.\n"
        "L'un sans l'autre ne fonctionne pas.",
        sz=10, c=LTGRAY)

    # ── LEVIER 2 : GOOGLE ADS ──
    x2 = x1 + cw + gap
    _card(sl, x2, cy, cw, ch2, CARD)
    _rect(sl, x2, cy, cw, Inches(.32), VIOLET)
    _tb(sl, x2, cy+Inches(.04), cw, Inches(.25),
        "DÈS QUE BUDGET DISPONIBLE", sz=10, c=WHITE, b=True, al=PP_ALIGN.CENTER)

    _tb(sl, x2+Inches(.2), cy+Inches(.45), cw-Inches(.4), Inches(.25),
        "Google Ads + Meta", sz=16, c=WHITE, b=True)

    _tb(sl, x2+Inches(.2), cy+Inches(.8), Inches(1.6), Inches(.2),
        "Budget", sz=9, c=GRAY3)
    _tb(sl, x2+Inches(.2), cy+Inches(1.0), Inches(2.5), Inches(.3),
        "dès 2 K€/mois", sz=24, c=VIOLET, b=True)
    _tb(sl, x2+Inches(.2), cy+Inches(1.35), cw-Inches(.4), Inches(.2),
        "scalable progressivement", sz=10, c=LTGRAY)

    _tb(sl, x2+Inches(.2), cy+Inches(1.65), Inches(1.6), Inches(.2),
        "Retour estimé", sz=9, c=GRAY3)
    _tb(sl, x2+Inches(.2), cy+Inches(1.85), Inches(2.5), Inches(.3),
        "×3 à 4 le budget", sz=20, c=VIOLET, b=True)
    _tb(sl, x2+Inches(.2), cy+Inches(2.2), cw-Inches(.4), Inches(.4),
        "2 K€ → ~7 K€ CA/mois\n"
        "5 K€ → ~18 K€ CA/mois",
        sz=11, c=WHITE)

    _tb(sl, x2+Inches(.2), cy+Inches(2.75), cw-Inches(.4), Inches(.2),
        "ROI", sz=9, c=GRAY3)
    _tb(sl, x2+Inches(.2), cy+Inches(2.9), Inches(2), Inches(.35),
        "×3 – 4", sz=20, c=VIOLET, b=True)

    _tb(sl, x2+Inches(.2), cy+Inches(3.35), cw-Inches(.4), Inches(.6),
        "Search + Shopping + saisonnalité Q4.\n"
        "Peut démarrer petit et scaler\n"
        "selon les résultats.",
        sz=10, c=LTGRAY)

    # ── LEVIER 3 : CRM ──
    x3 = x2 + cw + gap
    _card(sl, x3, cy, cw, ch3, CARD)
    _rect(sl, x3, cy, cw, Inches(.32), GREEN)
    _tb(sl, x3, cy+Inches(.04), cw, Inches(.25),
        "QUICK WIN", sz=10, c=WHITE, b=True, al=PP_ALIGN.CENTER)

    _tb(sl, x3+Inches(.2), cy+Inches(.45), cw-Inches(.4), Inches(.25),
        "CRM / Email", sz=16, c=WHITE, b=True)

    _tb(sl, x3+Inches(.2), cy+Inches(.8), Inches(1.6), Inches(.2),
        "Budget", sz=9, c=GRAY3)
    _tb(sl, x3+Inches(.2), cy+Inches(1.0), Inches(2.5), Inches(.3),
        "300 – 800 €/mois", sz=24, c=GREEN, b=True)
    _tb(sl, x3+Inches(.2), cy+Inches(1.35), cw-Inches(.4), Inches(.2),
        "outil + setup", sz=10, c=LTGRAY)

    _tb(sl, x3+Inches(.2), cy+Inches(1.65), Inches(1.6), Inches(.2),
        "CA additionnel", sz=9, c=GRAY3)
    _tb(sl, x3+Inches(.2), cy+Inches(1.85), Inches(2.5), Inches(.3),
        "+20 – 40 K€/an", sz=20, c=GREEN, b=True)
    _tb(sl, x3+Inches(.2), cy+Inches(2.2), cw-Inches(.4), Inches(.4),
        "Abandon panier : +5-10 % CA\n"
        "Relance Q4 CE/CSE",
        sz=11, c=WHITE)

    _tb(sl, x3+Inches(.2), cy+Inches(2.75), cw-Inches(.4), Inches(.2),
        "ROI", sz=9, c=GRAY3)
    _tb(sl, x3+Inches(.2), cy+Inches(2.9), Inches(2), Inches(.35),
        "×4 – 5", sz=20, c=GREEN, b=True)

    _tb(sl, x3+Inches(.2), cy+Inches(3.35), cw-Inches(.4), Inches(.6),
        "Welcome, abandon panier,\n"
        "post-achat, relance saisonnière.\n"
        "Rentable dès le mois 1.",
        sz=10, c=LTGRAY)

    # ── Bottom insight ──
    _tb(sl, Inches(.8), Inches(6.55), Inches(11.5), Inches(.3),
        "La refonte pose les fondations SEO. L'accompagnement mensuel génère la croissance. "
        "Les Ads et le CRM accélèrent le retour.",
        sz=13, c=DARK, b=True)

    _footer(sl, "Projections à M12 · Panier moyen 50,76 € · CVR cible 2,5 % · ROAS SEA ×3-4 · Sources : GSC + WooCommerce 2025")


# ═══════════════════════════════════════════════════
#  CORE SLIDE 9 — ROI GLOBAL DU PROJET
# ═══════════════════════════════════════════════════

def slide_roi(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    _title(sl, "Le projet en chiffres.", sz=36)

    # ── Top: Key assumptions (KPI strip) ──
    assumptions = [
        ("Panier moyen",   "50,76 €"),
        ("CVR actuel",     "~1,3 %"),
        ("CVR cible",      "2,5 %"),
        ("ROAS Ads",       "×3 – 4"),
    ]
    ax = Inches(.8)
    for label, val in assumptions:
        _card(sl, ax, Inches(1.4), Inches(2.4), Inches(.6), CARD_LT)
        _tb(sl, ax+Inches(.15), Inches(1.43), Inches(2.1), Inches(.2),
            label, sz=8, c=GRAY3)
        _tb(sl, ax+Inches(.15), Inches(1.62), Inches(2.1), Inches(.3),
            val, sz=16, c=DARK, b=True)
        ax += Inches(2.65)

    # ── Center: Waterfall progression ──
    wy = Inches(2.35)
    _tb(sl, Inches(.8), wy, Inches(6), Inches(.25),
        "PROGRESSION DU CA WEB — VISION À 12 MOIS", sz=9, c=GRAY3, b=True)

    # Waterfall data: (label, delta, cumul, color, detail)
    # Math: 7K×1.3%×50=4.6K → 7K×2.5%=8.75K → 9.5K×2.5%=11.9K → +3.6K×2.5%=+4.5K → +12%CRM
    waterfall = [
        ("Aujourd'hui",        55,   55,  GRAY3,  "7 000 sessions · CVR 1,3 % · 90 cmd/mois · +27 %/an"),
        ("Refonte (CVR ×2)",   50,  105,  DARK,   "Même trafic, meilleure conversion → 2,5 % CVR"),
        ("SEO Run (trafic)",   35,  140,  RED,    "+2 500 visites hors-marque/mois · contenu expert"),
        ("Ads (2K€/mois)",     55,  195,  VIOLET, "3 600 clics/mois · CPC 0,55 € · CVR 2,5 %"),
        ("CRM",                15,  210,  GREEN,  "Abandon panier + relance + post-achat · +12 % CA org."),
    ]

    max_val = 280  # scale reference
    bx_start = Inches(3.0)
    bw_max = Inches(9.2)
    row_h = Inches(.7)
    bar_h = Inches(.28)
    ry = wy + Inches(.35)

    cumul_prev = 0
    for label, delta, cumul, color, detail in waterfall:
        is_base = (cumul_prev == 0)
        y = ry

        # Label
        _tb(sl, Inches(.8), y+Inches(.05), Inches(2.1), Inches(.25),
            label, sz=12, c=DARK, b=True)

        # Cumulative bar (gray background)
        if not is_base:
            # Previous portion (gray)
            prev_w = bw_max * (cumul_prev / max_val)
            _bar(sl, bx_start, y+Inches(.08), prev_w, bar_h, BAR_BG)
            # Delta portion (colored)
            delta_w = bw_max * (delta / max_val)
            _bar(sl, bx_start + prev_w - Inches(.02), y+Inches(.08),
                 delta_w + Inches(.02), bar_h, color)
            # Delta label
            _tb(sl, bx_start + prev_w + delta_w + Inches(.1), y+Inches(.05),
                Inches(1.2), Inches(.25),
                f"+{delta} K€", sz=11, c=color, b=True)
        else:
            # Base bar
            base_w = bw_max * (cumul / max_val)
            _bar(sl, bx_start, y+Inches(.08), base_w, bar_h, color)
            _tb(sl, bx_start + base_w + Inches(.1), y+Inches(.05),
                Inches(1), Inches(.25),
                f"{cumul} K€", sz=11, c=color, b=True)

        # Detail text
        _tb(sl, Inches(.8), y+Inches(.32), Inches(2.1), Inches(.2),
            detail, sz=8, c=GRAY3)

        cumul_prev = cumul
        ry += row_h

    # ── Total bar ──
    ty = ry + Inches(.1)
    _hline(sl, Inches(.8), ty, Inches(11.5), RED)
    total_w = bw_max * (240 / max_val)
    _bar(sl, bx_start, ty+Inches(.12), total_w, Inches(.35), RED)

    _tb(sl, Inches(.8), ty+Inches(.1), Inches(2.1), Inches(.35),
        "OBJECTIF M12", sz=12, c=DARK, b=True)
    _tb(sl, bx_start + total_w + Inches(.15), ty+Inches(.1), Inches(2.5), Inches(.35),
        "~210 K€/an", sz=22, c=RED, b=True)

    # Multiplier
    _tb(sl, Inches(.8), ty+Inches(.5), Inches(11.5), Inches(.35),
        "×4 vs aujourd'hui — budget media 2 K€/mois (CPC 0,55 €). "
        "Projection prudente basée sur vos données réelles GSC + WooCommerce.",
        sz=14, c=GRAY2)

    _footer(sl, "Projections conservatrices · Panier 50,76 € · CVR cible 2,5 % · ROAS ×3-4 · Sources : GSC + WooCommerce 2025")


# ═══════════════════════════════════════════════════
#  CORE SLIDE 10 — PORTAIL DE MARQUE
# ═══════════════════════════════════════════════════

def slide_portail(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    _title(sl, "1 marque. 2 univers. 1 portail.", sz=40)
    _tb(sl, Inches(.8), Inches(1.25), Inches(11.5), Inches(.4),
        "64 900 internautes cherchent « la mère poulard » chaque mois.\n"
        "Ils doivent atterrir au bon endroit — sans perdre l'autorité SEO.",
        sz=14, c=GRAY2)

    # ── Flow diagram ──
    # Central hub
    hub_w, hub_h = Inches(3.2), Inches(1.3)
    hub_x = (SLIDE_W - hub_w) // 2
    hub_y = Inches(2.35)
    _card(sl, hub_x, hub_y, hub_w, hub_h, CARD)
    _tb(sl, hub_x, hub_y+Inches(.2), hub_w, Inches(.3),
        "lamerepoulard.fr", sz=20, c=WHITE, b=True, al=PP_ALIGN.CENTER)
    _tb(sl, hub_x, hub_y+Inches(.6), hub_w, Inches(.25),
        "64 900 recherches/mois", sz=12, c=RED, b=True, al=PP_ALIGN.CENTER)
    _tb(sl, hub_x, hub_y+Inches(.9), hub_w, Inches(.2),
        "Hub de marque · portail d'orientation", sz=9, c=GRAY3, al=PP_ALIGN.CENTER)

    # Arrow indicator (simple chevron text)
    arrow_y = hub_y + hub_h + Inches(.15)
    _tb(sl, Inches(3.5), arrow_y, Inches(2), Inches(.25),
        "▼", sz=16, c=RED, al=PP_ALIGN.CENTER)
    _tb(sl, Inches(7.5), arrow_y, Inches(2), Inches(.25),
        "▼", sz=16, c=MAGENTA, al=PP_ALIGN.CENTER)

    # Left destination: Biscuiterie
    dest_w = Inches(4.2)
    dest_h = Inches(1.4)
    dest_y = arrow_y + Inches(.35)
    dest_x1 = Inches(1.2)
    _card(sl, dest_x1, dest_y, dest_w, dest_h, CARD)
    _rect(sl, dest_x1, dest_y, dest_w, Inches(.28), RED)
    _tb(sl, dest_x1, dest_y+Inches(.03), dest_w, Inches(.22),
        "BISCUITERIE", sz=10, c=WHITE, b=True, al=PP_ALIGN.CENTER)
    _tb(sl, dest_x1+Inches(.2), dest_y+Inches(.4), dest_w-Inches(.4), Inches(.25),
        "biscuiterie-mere-poulard.com", sz=13, c=WHITE, b=True)
    _tb(sl, dest_x1+Inches(.2), dest_y+Inches(.7), dest_w-Inches(.4), Inches(.5),
        "E-commerce B2C · Coffrets & biscuits\nCE/CSE · LMP Professionnel\n50K€ CA web → objectif ×7",
        sz=10, c=LTGRAY)

    # Right destination: Auberge
    dest_x2 = Inches(7.7)
    _card(sl, dest_x2, dest_y, dest_w, dest_h, CARD)
    _rect(sl, dest_x2, dest_y, dest_w, Inches(.28), MAGENTA)
    _tb(sl, dest_x2, dest_y+Inches(.03), dest_w, Inches(.22),
        "AUBERGE & HÔTELLERIE", sz=10, c=WHITE, b=True, al=PP_ALIGN.CENTER)
    _tb(sl, dest_x2+Inches(.2), dest_y+Inches(.4), dest_w-Inches(.4), Inches(.25),
        "lamerepoulard.com", sz=13, c=WHITE, b=True)
    _tb(sl, dest_x2+Inches(.2), dest_y+Inches(.7), dest_w-Inches(.4), Inches(.5),
        "Hôtel · Restaurant · Expérience\nRéservation · Mont Saint-Michel",
        sz=10, c=LTGRAY)

    # ── 4 benefit cards at bottom ──
    by = Inches(5.3)
    bw = Inches(2.55)
    bh = Inches(1.15)
    bgap = Inches(.2)

    benefits = [
        ("Autorité SEO", RED,
         "L'autorité du domaine reste\nconcentrée. Pas de perte\nau moment de la migration."),
        ("GEO / IA ready", VIOLET,
         "1 entité = 1 réponse IA.\nLes LLMs citent la source\nla plus complète."),
        ("Hub scalable", GREEN,
         "LMP Pro, international,\nnouvelles marques — le portail\nabsorbe chaque extension."),
        ("Expérience claire", MAGENTA,
         "L'internaute choisit\nson univers en 1 clic.\nZéro confusion."),
    ]
    for i, (label, accent, desc) in enumerate(benefits):
        x = Inches(.8) + i * (bw + bgap)
        _card(sl, x, by, bw, bh, CARD)
        _tb(sl, x+Inches(.15), by+Inches(.1), bw-Inches(.3), Inches(.22),
            label, sz=11, c=accent, b=True)
        _tb(sl, x+Inches(.15), by+Inches(.38), bw-Inches(.3), Inches(.7),
            desc, sz=9, c=LTGRAY)

    _footer(sl, "Architecture domaine · Recommandation SLASHR + Fractory")


# ═══════════════════════════════════════════════════
#  CORE SLIDE 11 — MÉTHODOLOGIE
# ═══════════════════════════════════════════════════

def slide_methodologie(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    _title(sl, "Comment ces chiffres sont calculés.", sz=34)

    # ── The formula ──
    fy = Inches(1.5)
    _card(sl, Inches(.8), fy, Inches(11.5), Inches(.6), CARD)
    _tb(sl, Inches(1.1), fy+Inches(.12), Inches(10.5), Inches(.35),
        "Trafic mensuel   ×   Taux de conversion   ×   Panier moyen   =   CA mensuel",
        sz=18, c=WHITE, b=True, al=PP_ALIGN.CENTER)

    # ── Calculation per lever ──
    ty = Inches(2.4)
    _tb(sl, Inches(.8), ty, Inches(6), Inches(.25),
        "DÉTAIL PAR LEVIER", sz=9, c=GRAY3, b=True)

    # Table header
    cols = [Inches(1.6), Inches(2.4), Inches(1.5), Inches(1.0), Inches(1.0), Inches(2.0), Inches(1.2)]
    headers = ["Levier", "Objectif", "Trafic", "CVR", "Panier", "Calcul", "CA estimé"]
    hx = Inches(.5)
    hy = ty + Inches(.3)
    for j, hdr in enumerate(headers):
        _tb(sl, hx + sum(cols[:j]), hy, cols[j], Inches(.25),
            hdr, sz=9, c=GRAY3, b=True)
    _hline(sl, hx, hy + Inches(.28), sum(cols), RED)

    # Table rows
    # (label, objectif, trafic, cvr, panier, calc, ca, accent)
    rows = [
        ("Aujourd'hui",
         "GSC + GA4 + WooCommerce",
         "7 000", "1,3 %", "50,76 €",
         "7 000 × 1,3 % × 50 €", "~4,6 K€",
         GRAY3),
        ("+ Refonte",
         "CVR ×2 → 2,5 %\nBenchmark food : 4,6-6,2 %",
         "7 000", "2,5 %", "50 €",
         "7 000 × 2,5 % × 50 €", "~8,8 K€",
         DARK),
        ("+ SEO Run",
         "+2,5K visites hors-marque\nAuj. <100 clics HM (GSC)",
         "9 500", "2,5 %", "50 €",
         "9 500 × 2,5 % × 50 €", "~11,9 K€",
         RED),
        ("+ Ads",
         "Budget 2K€/mois\nCPC 0,55 € · 3 600 clics",
         "+3 600", "2,5 %", "50 €",
         "3 600 × 2,5 % × 50 €", "+4,5 K€",
         VIOLET),
        ("+ CRM",
         "+12 % CA en relance\nBenchmark : 20-30 % (Klaviyo)",
         "—", "+12 %", "—",
         "12 % × CA organique", "+1,4 K€",
         GREEN),
    ]

    rh = Inches(.46)
    ry = hy + Inches(.35)
    for label, objectif, trafic, cvr, panier, calc, ca, accent in rows:
        vals = [label, objectif, trafic, cvr, panier, calc, ca]
        for j, val in enumerate(vals):
            is_first = (j == 0)
            is_obj = (j == 1)
            is_last = (j == len(vals) - 1)
            _tb(sl, hx + sum(cols[:j]), ry, cols[j], rh,
                val,
                sz=10 if is_first else (9 if is_obj else 9),
                c=accent if (is_first or is_last or is_obj) else DARK,
                b=(is_first or is_last))
        ry += rh

    # Total row
    _hline(sl, hx, ry + Inches(.02), sum(cols), RED)
    ry += Inches(.1)
    _tb(sl, hx, ry, cols[0], rh,
        "TOTAL M12", sz=11, c=RED, b=True)
    _tb(sl, hx + cols[0], ry, cols[1], rh,
        "×4 vs aujourd'hui", sz=10, c=RED, b=True)
    _tb(sl, hx + sum(cols[:5]), ry, cols[5], rh,
        "~17,5 K€/mois", sz=11, c=RED, b=True)
    _tb(sl, hx + sum(cols[:6]), ry, cols[6], rh,
        "~210 K€", sz=11, c=RED, b=True)

    # ── Where do these numbers come from ──
    sy = Inches(5.3)
    _tb(sl, Inches(.8), sy, Inches(6), Inches(.25),
        "D'OÙ VIENNENT CES DONNÉES", sz=9, c=GRAY3, b=True)

    sources = [
        ("Trafic : 7 000/mois",  "GSC + GA4",             "4 500 org. (GSC réel)\n+2 500 direct/ref (GA4 ratio)"),
        ("CVR actuel : 1,3 %",   "WooCommerce + GA4",     "1 082 cmd / ~84K sessions\nCible 2,5 % (food : 4,6-6,2 %)"),
        ("SEO : +2,5K visites",  "DataForSEO + GSC",      "Auj. <100 clics HM (GSC)\nProjection prudente M12"),
        ("Ads : 2K€ · 3,6K clics", "CPC 0,55 €",          "ROAS ×2,3 (55K CA / 24K budget)\nCVR 2,5 % · panier 50 €"),
        ("CRM : +12 % CA",       "Klaviyo 2025",          "Benchmark email : 20-30 %\nAn 1 depuis zéro = 12 %"),
    ]
    sw = Inches(2.16)
    for i, (what, source, detail) in enumerate(sources):
        x = Inches(.8) + i * (sw + Inches(.1))
        _card(sl, x, sy + Inches(.3), sw, Inches(.85), CARD)
        _tb(sl, x + Inches(.12), sy + Inches(.35), sw - Inches(.24), Inches(.2),
            what, sz=9, c=GRAY3)
        _tb(sl, x + Inches(.12), sy + Inches(.55), sw - Inches(.24), Inches(.25),
            source, sz=11, c=WHITE, b=True)
        _tb(sl, x + Inches(.12), sy + Inches(.8), sw - Inches(.24), Inches(.3),
            detail, sz=8, c=LTGRAY)

    _footer(sl, "Données réelles : GSC 12 mois + WooCommerce 2024-2025 + GA4 · Projections mesurables via GA4 + GSC + Google Ads + CRM")


# ═══════════════════════════════════════════════════
#  ANNEXE — SÉPARATEUR
# ═══════════════════════════════════════════════════

def annexe_separator(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    _tb(sl, Inches(.8), Inches(2.5), Inches(11.5), Inches(1),
        "Annexes", sz=52, c=DARK, b=True, al=PP_ALIGN.CENTER)
    _tb(sl, Inches(.8), Inches(3.7), Inches(11.5), Inches(.5),
        "Données détaillées · Benchmark · Territoires · Diagnostic",
        sz=16, c=GRAY3, al=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
#  ANNEXE A1 — BENCHMARK CONCURRENTIEL
# ═══════════════════════════════════════════════════

def annexe_benchmark(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    _annexe_label(sl, 1, "BENCHMARK CONCURRENTIEL")
    _title(sl, "Trafic organique : marque vs hors-marque", sz=30, y=.45)

    competitors = [
        ("La Trinitaine",          20016, .82),
        ("Biscuiterie de Kerlann",  6126, .17),
        ("La Mère Poulard",         4500, .02),
        ("Galettes Penven",          2287, .79),
    ]
    max_v = 20016
    bl, bw_max = Inches(3.2), Inches(8.0)
    bar_h, rh = Inches(.35), Inches(.6)
    bt = Inches(2.0)

    for i, (name, total, pct_nb) in enumerate(competitors):
        y = bt + rh * i
        is_lmp = "Poulard" in name
        _tb(sl, Inches(.8), y, Inches(2.3), Inches(.3),
            name, sz=12, c=DARK if is_lmp else GRAY2, b=is_lmp)
        full_w = bw_max * (total / max_v)
        brand_w = full_w * (1 - pct_nb)
        nb_w = full_w * pct_nb
        if brand_w > Inches(.05):
            _bar(sl, bl, y+Inches(.02), brand_w, bar_h, DARK)
        if nb_w > Inches(.05):
            _bar(sl, bl+brand_w-Inches(.02), y+Inches(.02), nb_w+Inches(.02), bar_h, RED)
        _tb(sl, bl+full_w+Inches(.15), y, Inches(3), Inches(.3),
            f"{total:,} · {int(pct_nb*100)}% hors-marque".replace(",", " "),
            sz=11, c=DARK if is_lmp else GRAY2, b=is_lmp)

    # Legend
    ly = bt + rh * 4 + Inches(.15)
    _rect(sl, Inches(3.2), ly, Inches(.3), Inches(.15), DARK)
    _tb(sl, Inches(3.6), ly-Inches(.02), Inches(1.2), Inches(.2), "Marque", sz=9, c=GRAY3)
    _rect(sl, Inches(4.8), ly, Inches(.3), Inches(.15), RED)
    _tb(sl, Inches(5.2), ly-Inches(.02), Inches(1.5), Inches(.2), "Hors-marque", sz=9, c=GRAY3)

    _bullet(sl, Inches(1.1), Inches(4.9), Inches(10.5),
            "La Trinitaine : ~16 000 clics hors-marque/mois (DataForSEO). "
            "La Mère Poulard : <100 (GSC réel). Gap ×100+.", sz=13, c=GRAY2)
    _bullet(sl, Inches(1.1), Inches(5.45), Inches(10.5),
            "La notoriété de marque LMP est forte (64 900 rech./mois), "
            "mais elle ne rayonne pas au-delà de la marque.", sz=13, c=GRAY2)

    _footer(sl, "Sources : GSC 12 mois (LMP) · DataForSEO (concurrents) · 2025")


# ═══════════════════════════════════════════════════
#  ANNEXE A2 — TERRITOIRES SEO
# ═══════════════════════════════════════════════════

def annexe_territoires(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    _annexe_label(sl, 2, "TERRITOIRES SEO")
    _title(sl, "6 territoires d'achat à conquérir", sz=30, y=.45)

    territories = [
        ("biscuits",         "40 500", "#3"),
        ("biscuiterie",      "14 800", "—"),
        ("petit beurre",      "8 100", "#9"),
        ("panier gourmand",   "5 400", "—"),
        ("coffret gourmand",  "2 900", "—"),
        ("biscuiterie MSM",   "1 000", "—"),
    ]
    tx, ty = Inches(.8), Inches(1.7)
    cols = [Inches(2.4), Inches(1.6), Inches(1.6)]

    for j, hdr in enumerate(["Territoire", "Volume/mois", "Position LMP"]):
        _tb(sl, tx+sum(cols[:j]), ty, cols[j], Inches(.3), hdr, sz=10, c=GRAY3, b=True)
    _hline(sl, tx, ty+Inches(.32), sum(cols), RED)

    rh = Inches(.36)
    for i, (term, vol, pos) in enumerate(territories):
        y = ty + Inches(.4) + rh * i
        _tb(sl, tx, y, cols[0], rh, term, sz=13, c=DARK, b=True)
        _tb(sl, tx+cols[0], y, cols[1], rh, vol, sz=13, c=RED, b=True)
        pc = GREEN if pos != "—" else GRAY3
        _tb(sl, tx+cols[0]+cols[1], y, cols[2], rh, pos, sz=13, c=pc, b=(pos != "—"))

    # Shopping callout
    cx, cy2, cw2, ch2 = Inches(7), Inches(1.7), Inches(5.3), Inches(2.4)
    _card(sl, cx, cy2, cw2, ch2, CARD)
    _tb(sl, cx+Inches(.25), cy2+Inches(.15), cw2-Inches(.5), Inches(.25),
        "GOOGLE SHOPPING", sz=10, c=MAGENTA, b=True)
    _tb(sl, cx+Inches(.25), cy2+Inches(.5), cw2-Inches(.5), Inches(.7),
        "Vos coffrets apparaissent dans Shopping.\nMais c'est La Grande Épicerie qui vend.",
        sz=13, c=WHITE)
    _tb(sl, cx+Inches(.25), cy2+Inches(1.4), cw2-Inches(.5), Inches(.7),
        "Sans feed produit optimisé (prix, avis, disponibilité),\n"
        "les clics transactionnels vont aux revendeurs.",
        sz=10, c=LTGRAY)

    # Product mapping
    my = Inches(4.6)
    _tb(sl, Inches(.8), my, Inches(11.5), Inches(.35),
        "Chemin Search par gamme produit", sz=14, c=DARK, b=True)
    gammes = [
        ("Biscuits & Sablés",  "SEO + Shopping",        RED),
        ("Coffrets & Paniers", "SEO + SEA saisonnier",  MAGENTA),
        ("Madeleines",         "Contenu recette → CTA", VIOLET),
        ("Gamme Normandie",    "SEO territorial",       GREEN),
    ]
    gw, gh = Inches(2.7), Inches(.95)
    for i, (gamme, chemin, accent) in enumerate(gammes):
        x = Inches(.8) + i * (gw + Inches(.17))
        _card(sl, x, my+Inches(.45), gw, gh, CARD)
        _tb(sl, x+Inches(.15), my+Inches(.55), gw-Inches(.3), Inches(.3),
            gamme, sz=12, c=WHITE, b=True)
        _tb(sl, x+Inches(.15), my+Inches(.9), gw-Inches(.3), Inches(.3),
            chemin, sz=10, c=accent, b=True)

    _footer(sl, "Source : DataForSEO · SERP analysis 6 keywords commerciaux · mars 2026")


# ═══════════════════════════════════════════════════
#  ANNEXE A3 — CONTENU RECETTES
# ═══════════════════════════════════════════════════

def annexe_contenu(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    _annexe_label(sl, 3, "CONTENU INFORMATIONNEL")
    _title(sl, "67K internautes cherchent vos recettes", sz=30, y=.45)

    info_kws = [
        ("recette madeleine",   49500),
        ("recette galette",      9900),
        ("recette sablé",        4400),
        ("spécialité normande",  3600),
    ]
    max_v = 49500
    bt, lw, bl, bmw = Inches(1.8), Inches(2.5), Inches(3.2), Inches(7)
    rh = Inches(.55)
    for i, (kw, vol) in enumerate(info_kws):
        y = bt + rh * i
        _tb(sl, Inches(.8), y, lw, Inches(.3), kw, sz=13, c=DARK)
        w = bmw * (vol / max_v)
        _bar(sl, bl, y+Inches(.03), w, Inches(.24), RED)
        _tb(sl, bl+w+Inches(.12), y, Inches(1.5), Inches(.3),
            f"{vol:,}/mois".replace(",", " "), sz=10, c=GRAY3)

    # Funnel
    fy = Inches(4.2)
    _tb(sl, Inches(.8), fy-Inches(.35), Inches(6), Inches(.3),
        "Stratégie de contenu → conversion", sz=14, c=DARK, b=True)
    steps = [
        ("Requête info",   '"recette sablé\npur beurre"',            GRAY3),
        ("Contenu expert", "Guide + recette\nE-E-A-T (depuis 1888)", VIOLET),
        ("CTA produit",    '"Goûtez\nl\'original"',                  RED),
        ("Conversion",     "Achat\ne-commerce",                      GREEN),
    ]
    sw, sh, aw = Inches(2.5), Inches(.95), Inches(.35)
    x = Inches(.8)
    for j, (label, desc, accent) in enumerate(steps):
        if j > 0:
            _tb(sl, x, fy+Inches(.15), aw, Inches(.4), "→", sz=22, c=GRAY3, al=PP_ALIGN.CENTER)
            x += aw
        _card(sl, x, fy, sw, sh, CARD)
        _tb(sl, x+Inches(.1), fy+Inches(.08), sw-Inches(.2), Inches(.22),
            label, sz=10, c=accent, b=True, al=PP_ALIGN.CENTER)
        _tb(sl, x+Inches(.1), fy+Inches(.32), sw-Inches(.2), Inches(.55),
            desc, sz=9, c=LTGRAY, al=PP_ALIGN.CENTER)
        x += sw + Inches(.08)

    _tb(sl, Inches(.8), Inches(5.6), Inches(11.5), Inches(.5),
        "E-E-A-T natif : marque centenaire (1888), savoir-faire documenté.\n"
        "0 concurrent artisanal positionné sur ces requêtes.",
        sz=14, c=GRAY2)
    _footer(sl, "Source : DataForSEO · search_intent · mars 2026")


# ═══════════════════════════════════════════════════
#  ANNEXE A4 — ENJEU BRETON
# ═══════════════════════════════════════════════════

def annexe_breton(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    _annexe_label(sl, 4, "ENJEU BRETON")
    _title(sl, "69K recherches bretonnes :\ninaccessibles en SEO, captables en SEA", sz=28, y=.4)

    breton_kws = [
        ("palet breton",     22200),
        ("galette bretonne", 22200),
        ("gâteau breton",    18100),
        ("sablé breton",      5400),
    ]
    max_v = 22200
    bt, bl, bmw = Inches(2.0), Inches(3.2), Inches(4.2)
    rh = Inches(.5)
    for i, (kw, vol) in enumerate(breton_kws):
        y = bt + rh * i
        _tb(sl, Inches(.8), y, Inches(2.3), Inches(.3), kw, sz=13, c=DARK)
        w = bmw * (vol / max_v)
        _bar(sl, bl, y+Inches(.02), w, Inches(.22), RED)
        _tb(sl, bl+w+Inches(.12), y, Inches(1.5), Inches(.3),
            f"{vol:,}/mois".replace(",", " "), sz=10, c=GRAY3)

    sx, sw = Inches(8.3), Inches(4)
    _kpi(sl, sx, Inches(2.0), sw, Inches(.9),
         "CPC 0,14 – 0,52 €", "Coût par clic très bas", GREEN)
    _kpi(sl, sx, Inches(3.05), sw, Inches(.9),
         "~2 500 €/mois", "Budget Search Ads estimé", RED)
    _kpi(sl, sx, Inches(4.1), sw, Inches(.9),
         "~3 500 clics", "Clics qualifiés par mois", RED)

    sy = Inches(4.6)
    _card(sl, Inches(.8), sy, Inches(7), Inches(1.15), CARD)
    _tb(sl, Inches(1.1), sy+Inches(.1), Inches(6.4), Inches(.25),
        "SAISONNALITÉ", sz=10, c=RED, b=True)
    _tb(sl, Inches(1.1), sy+Inches(.4), Inches(6.4), Inches(.6),
        "Pics Noël & Pâques : coffrets +404 %, paniers +311 %\n"
        "→ Activer le SEA 6 semaines avant chaque pic",
        sz=12, c=WHITE)

    _tb(sl, Inches(.8), Inches(6.1), Inches(11.5), Inches(.4),
        "LMP est normande — « breton » ne peut pas apparaître sur le site. "
        "Le SEA cible ces requêtes sans contrainte éditoriale.",
        sz=11, c=GRAY2)
    _footer(sl, "Source : DataForSEO · search_intent + Google Ads CPC · mars 2026")


# ═══════════════════════════════════════════════════
#  ANNEXE A5 — INTERNATIONAL
# ═══════════════════════════════════════════════════

def annexe_international(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    _annexe_label(sl, 5, "INTERNATIONAL")
    _title(sl, "35 % de votre trafic est déjà international.", sz=30, y=.45)

    # GSC reality strip
    _card(sl, Inches(.8), Inches(1.5), Inches(11.5), Inches(.6), CARD)
    _tb(sl, Inches(1.1), Inches(1.55), Inches(2), Inches(.2),
        "DONNÉES GSC RÉELLES (12 MOIS)", sz=9, c=RED, b=True)
    _tb(sl, Inches(1.1), Inches(1.78), Inches(10.5), Inches(.25),
        "19 000 clics internationaux · Page /en/ = 2e page la plus visitée (7 305 clics, 255K impressions)",
        sz=12, c=WHITE)

    # Top countries
    kw, kh, ky = Inches(2.1), Inches(1.0), Inches(2.5)
    countries = [
        ("USA", "2 210 clics", "89K impr.", RED),
        ("UK", "1 950 clics", "35K impr.", RED),
        ("Canada", "1 335 clics", "17K impr.", VIOLET),
        ("Allemagne", "1 202 clics", "18K impr.", VIOLET),
        ("Italie", "977 clics", "9K impr.", GRAY3),
    ]
    for i, (country, clicks, impr, accent) in enumerate(countries):
        x = Inches(.8) + i * (kw + Inches(.12))
        _card(sl, x, ky, kw, kh, CARD)
        _tb(sl, x+Inches(.12), ky+Inches(.1), kw-Inches(.24), Inches(.2),
            country, sz=11, c=accent, b=True)
        _tb(sl, x+Inches(.12), ky+Inches(.35), kw-Inches(.24), Inches(.2),
            clicks, sz=14, c=WHITE, b=True)
        _tb(sl, x+Inches(.12), ky+Inches(.6), kw-Inches(.24), Inches(.2),
            impr, sz=10, c=LTGRAY)

    ay = Inches(4.0)
    advantages = [
        ('Pas de contrainte "breton"',
         "En anglais, pas de conflit territorial.\nPositionnement SEO + SEA libre.", GREEN),
        ("CPC ultra-bas",
         "0,12 – 0,55 € par clic.\nMarché peu concurrentiel.", RED),
        ("Déjà 23K recherches EN/mois",
         "Marque + catégorie.\nLe trafic existe, le contenu EN est mince.", VIOLET),
    ]
    aw, ah = Inches(3.6), Inches(1.2)
    for i, (title, desc, accent) in enumerate(advantages):
        x = Inches(.8) + i * (aw + Inches(.17))
        _card(sl, x, ay, aw, ah, CARD)
        _tb(sl, x+Inches(.2), ay+Inches(.12), aw-Inches(.4), Inches(.4),
            title, sz=13, c=accent, b=True)
        _tb(sl, x+Inches(.2), ay+Inches(.55), aw-Inches(.4), Inches(.55),
            desc, sz=11, c=LTGRAY)

    _footer(sl, "Sources : GSC 12 mois (clics réels) + DataForSEO (volumes EN) · 2025")


# ═══════════════════════════════════════════════════
#  ANNEXE A6 — B2B
# ═══════════════════════════════════════════════════

def annexe_b2b(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    _annexe_label(sl, 6, "AUDIENCES B2B")
    _title(sl, "CE/CSE et Professionnels : poser les fondations", sz=28, y=.4)

    col_w = Inches(5.5)
    cy = Inches(1.7)

    # CE/CSE
    col_l = Inches(.8)
    _card(sl, col_l, cy, col_w, Inches(3.5), CARD)
    _tb(sl, col_l+Inches(.25), cy+Inches(.12), col_w-Inches(.5), Inches(.25),
        "CE / CSE & ASSOCIATIONS", sz=10, c=RED, b=True)
    _tb(sl, col_l+Inches(.25), cy+Inches(.42), col_w-Inches(.5), Inches(.5),
        "Commandes groupées, coffrets cadeaux collaborateurs.\n"
        "Aujourd'hui : 90 % des commandes par email,\n"
        "ressaisie manuelle dans l'ERP.", sz=11, c=LTGRAY)
    cse_kws = [
        ("coffret cadeau CSE",       "110/mois",  "CPC 4,54 €"),
        ("coffret cadeau entreprise", "170/mois",  "CPC 2,10 €"),
        ("cadeau gastronomique",      "140/mois",  "CPC 1,76 €"),
    ]
    qy = cy + Inches(1.35)
    _tb(sl, col_l+Inches(.25), qy, Inches(3), Inches(.2),
        "REQUÊTES TYPIQUES", sz=9, c=GRAY3, b=True)
    for i, (kw, vol, cpc) in enumerate(cse_kws):
        y = qy + Inches(.3) + Inches(.3) * i
        _tb(sl, col_l+Inches(.25), y, Inches(2.8), Inches(.25), kw, sz=11, c=WHITE)
        _tb(sl, col_l+Inches(3.1), y, Inches(1), Inches(.25), vol, sz=11, c=RED, b=True)
        _tb(sl, col_l+Inches(4.2), y, Inches(1), Inches(.25), cpc, sz=9, c=GRAY3)
    _tb(sl, col_l+Inches(.25), cy+Inches(2.7), col_w-Inches(.5), Inches(.55),
        "Ultra-qualifié : CPC élevé = forte intention d'achat.\n"
        "1 page CE/CSE indexable capte ce trafic sans effort continu.",
        sz=11, c=GREEN)

    # Pro B2B
    col_r = Inches(6.55)
    _card(sl, col_r, cy, col_w, Inches(3.5), CARD)
    _tb(sl, col_r+Inches(.25), cy+Inches(.12), col_w-Inches(.5), Inches(.25),
        "PROFESSIONNELS · CHR · FOODSERVICE", sz=10, c=VIOLET, b=True)
    _tb(sl, col_r+Inches(.25), cy+Inches(.42), col_w-Inches(.5), Inches(.5),
        "Traiteurs, restaurateurs, grossistes.\n"
        "Intéressés par Dutoit (prêt-à-garnir)\n"
        "et Tourniayre (cornets) — marques B2B du groupe.", sz=11, c=LTGRAY)
    pro_kws = [
        ("grossiste pâtisserie", "1 300/mois", ""),
        ("grossiste biscuit",      "320/mois", ""),
        ("cornet de glace pro",     "70/mois", "Dutoit"),
        ("fond de tarte pro",       "50/mois", "Tourniayre"),
    ]
    _tb(sl, col_r+Inches(.25), qy, Inches(3), Inches(.2),
        "REQUÊTES TYPIQUES", sz=9, c=GRAY3, b=True)
    for i, (kw, vol, note) in enumerate(pro_kws):
        y = qy + Inches(.3) + Inches(.3) * i
        _tb(sl, col_r+Inches(.25), y, Inches(2.8), Inches(.25), kw, sz=11, c=WHITE)
        _tb(sl, col_r+Inches(3.1), y, Inches(1), Inches(.25), vol, sz=11, c=VIOLET, b=True)
        if note:
            _tb(sl, col_r+Inches(4.2), y, Inches(1), Inches(.25), note, sz=9, c=GRAY3)
    _tb(sl, col_r+Inches(.25), cy+Inches(2.7), col_w-Inches(.5), Inches(.55),
        "Sans espace Dutoit/Tourniayre indexable,\n"
        "ces requêtes sont captées par d'autres fournisseurs.",
        sz=11, c=MAGENTA)

    _tb(sl, Inches(.8), Inches(5.6), Inches(11.5), Inches(.5),
        "Le B2B n'est pas la priorité 2026. Mais la refonte est le moment de poser "
        "les fondations : pages CE/CSE + espace pro indexables dès le lancement.",
        sz=13, c=GRAY2)
    _footer(sl, "Source : DataForSEO · search_intent + SERP analysis · mars 2026")


# ═══════════════════════════════════════════════════
#  ANNEXE A7 — GEO / LLMO
# ═══════════════════════════════════════════════════

def annexe_geo(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    _annexe_label(sl, 7, "GEO / LLMO")
    _title(sl, "Les IA répondent à la place de Google", sz=30, y=.45)

    _tb(sl, Inches(.8), Inches(1.5), Inches(11.5), Inches(.5),
        "ChatGPT, Perplexity, Google AI Overviews : quand un internaute pose une question,\n"
        "l'IA cite des marques. Si vous n'y êtes pas, un concurrent y sera.",
        sz=13, c=GRAY2)

    platforms = [
        ("ChatGPT",             "740 M",  "utilisateurs/mois",
         "Recommande des marques\ndans ses réponses", VIOLET),
        ("Google AI Overviews", "1 Md+",  "recherches avec IA/jour",
         "Résumé IA au-dessus\ndes résultats classiques", RED),
        ("Perplexity",          "100 M+", "requêtes/mois",
         "Réponses sourcées avec\nliens vers les sites cités", MAGENTA),
    ]
    pw, ph, py = Inches(3.6), Inches(1.7), Inches(2.4)
    for i, (name, val, unit, desc, accent) in enumerate(platforms):
        x = Inches(.8) + i * (pw + Inches(.17))
        _card(sl, x, py, pw, ph, CARD)
        _tb(sl, x+Inches(.2), py+Inches(.12), pw-Inches(.4), Inches(.25),
            name, sz=11, c=WHITE, b=True)
        _tb(sl, x+Inches(.2), py+Inches(.4), pw-Inches(.4), Inches(.4),
            val, sz=26, c=accent, b=True)
        _tb(sl, x+Inches(.2), py+Inches(.82), pw-Inches(.4), Inches(.2),
            unit, sz=9, c=GRAY3)
        _tb(sl, x+Inches(.2), py+Inches(1.1), pw-Inches(.4), Inches(.5),
            desc, sz=10, c=LTGRAY)

    ay = Inches(4.5)
    _card(sl, Inches(.8), ay, Inches(5.8), Inches(1.1), CARD)
    _tb(sl, Inches(1.1), ay+Inches(.08), Inches(5.3), Inches(.25),
        "L'AVANTAGE LA MÈRE POULARD", sz=10, c=GREEN, b=True)
    _tb(sl, Inches(1.1), ay+Inches(.35), Inches(5.3), Inches(.65),
        "Marque centenaire (1888), savoir-faire artisanal, Mont-Saint-Michel.\n"
        "Les LLMs valorisent l'E-E-A-T : expérience, expertise, autorité, fiabilité.\n"
        "Votre histoire est un actif GEO naturel — encore faut-il le structurer.",
        sz=11, c=WHITE)

    _card(sl, Inches(6.85), ay, Inches(5.5), Inches(1.1), CARD)
    _tb(sl, Inches(7.1), ay+Inches(.08), Inches(5), Inches(.25),
        "LE RISQUE", sz=10, c=MAGENTA, b=True)
    _tb(sl, Inches(7.1), ay+Inches(.35), Inches(5), Inches(.65),
        "Sans contenu structuré ni données Schema, les IA\n"
        "ne peuvent pas vous citer. Elles citent ceux qui publient\n"
        "du contenu expert : guides, recettes, fiches enrichies.",
        sz=11, c=WHITE)

    _footer(sl, "GEO = Generative Engine Optimization · Le SEO des réponses IA")


# ═══════════════════════════════════════════════════
#  ANNEXE A8 — DIAGNOSTIC S7
# ═══════════════════════════════════════════════════

def _radar_png():
    labels = ["S1\nIntentions", "S2\nArchitecture", "S3\nContenu",
              "S4\nUX/CVR", "S5\nAutorité", "S6\nDiffusion", "S7\nAmplification"]
    scores = [1, 2, 1, 2, 3, 1, 0]
    N = len(labels)
    angles = [n / N * 2 * math.pi for n in range(N)]
    angles += angles[:1]
    vals = scores + scores[:1]

    fig, ax = plt.subplots(figsize=(5.5, 5.5), subplot_kw=dict(polar=True))
    fig.patch.set_facecolor("#1A1A1A")
    ax.set_facecolor("#1A1A1A")
    ax.set_ylim(0, 5)
    ax.set_yticks([1, 2, 3, 4, 5])
    ax.set_yticklabels(["1", "2", "3", "4", "5"], color="#666", size=8)
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(labels, color="#B3B3B3", size=9, ha="center", fontfamily="sans-serif")
    ax.spines["polar"].set_color("#3A3D47")
    ax.grid(color="#3A3D47", linewidth=.5)
    ax.tick_params(axis="y", colors="#666")
    ax.plot(angles, vals, "o-", color="#E74601", linewidth=2.5, markersize=7)
    ax.fill(angles, vals, alpha=.12, color="#E74601")
    ax.plot([angles[4]], [3], "o", color="#22C55E", markersize=12, zorder=5)

    buf = io.BytesIO()
    plt.savefig(buf, format="png", dpi=180, bbox_inches="tight",
                facecolor="#1A1A1A", edgecolor="none")
    plt.close(); buf.seek(0)
    return buf


def annexe_diagnostic(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(sl)
    _annexe_label(sl, 8, "DIAGNOSTIC S7")
    _title(sl, "Diagnostic Search complet (S7)", sz=30, y=.45)

    radar = _radar_png()
    sl.shapes.add_picture(radar, Inches(5.5), Inches(1.3), Inches(7), Inches(5.5))

    cx, cw, ch = Inches(.8), Inches(5), Inches(1.35)

    _card(sl, cx, Inches(1.7), cw, ch, CARD_LT)
    _tb(sl, cx+Inches(.25), Inches(1.8), Inches(2), Inches(.22),
        "Actif", sz=10, c=GREEN, b=True)
    _tb(sl, cx+Inches(.25), Inches(2.05), cw-Inches(.5), Inches(.3),
        "Autorité (S5) : 3/5", sz=16, c=DARK, b=True)
    _tb(sl, cx+Inches(.25), Inches(2.4), cw-Inches(.5), Inches(.4),
        '#3 sur "biscuits" (40 500/mois). Fondation\n'
        "solide mais inexploitée.", sz=11, c=GRAY2)

    _card(sl, cx, Inches(3.25), cw, ch, CARD_LT)
    _tb(sl, cx+Inches(.25), Inches(3.35), Inches(2.5), Inches(.22),
        "Contrainte principale", sz=10, c=RED, b=True)
    _tb(sl, cx+Inches(.25), Inches(3.6), cw-Inches(.5), Inches(.3),
        "Contenu (S3) : 1/5", sz=16, c=DARK, b=True)
    _tb(sl, cx+Inches(.25), Inches(3.95), cw-Inches(.5), Inches(.4),
        "157 mots/page, 0 contenu ciblé par intention.\n"
        "Le catalogue est muet.", sz=11, c=GRAY2)

    _card(sl, cx, Inches(4.8), cw, ch, CARD_LT)
    _tb(sl, cx+Inches(.25), Inches(4.9), Inches(2), Inches(.22),
        "Leviers", sz=10, c=VIOLET, b=True)
    _tb(sl, cx+Inches(.25), Inches(5.15), cw-Inches(.5), Inches(.3),
        "Intentions (S1) + Architecture (S2)", sz=16, c=DARK, b=True)
    _tb(sl, cx+Inches(.25), Inches(5.5), cw-Inches(.5), Inches(.4),
        "Ciblage intentionnel + Schema Product.\n"
        "La refonte e-commerce est la fenêtre.", sz=11, c=GRAY2)

    _footer(sl, "Diagnostic S7 · DataForSEO + crawl · mars 2026")


# ═══════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── Core (12 slides) ──
    slide_geo_core(prs)       # 1. Le Search change
    slide_constat(prs)        # 2. 308K, 2%
    slide_pourquoi(prs)       # 3. Diagnostic
    slide_concurrence(prs)    # 4. Concurrence
    slide_identite(prs)       # 5. Identité Google : 3 piliers
    slide_seo(prs)            # 6. Approche SEO
    slide_sea(prs)            # 7. Approche SEA
    slide_synergie(prs)       # 8. Effet composé
    slide_projection(prs)     # 9. Où investir
    slide_roi(prs)            # 10. ROI global du projet
    slide_portail(prs)        # 11. Portail de marque
    slide_methodologie(prs)   # 12. Méthodologie

    # ── Annexe (separator + 8 slides) ──
    annexe_separator(prs)
    annexe_benchmark(prs)       # A1
    annexe_territoires(prs)     # A2
    annexe_contenu(prs)         # A3
    annexe_breton(prs)          # A4
    annexe_international(prs)   # A5
    annexe_b2b(prs)             # A6
    annexe_geo(prs)             # A7
    annexe_diagnostic(prs)      # A8

    out = "/Users/quentin/Downloads/SLASHR-SEO-SEA-LaMerePoulard.pptx"
    prs.save(out)
    print(f"Done — {out}")
    print("12 core + 1 separator + 8 annexes = 21 slides · 16:9 · fond blanc")


if __name__ == "__main__":
    main()
